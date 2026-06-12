import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def create_premium_presentation():
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Palette definition (Gris, Naranjo, Terracota, Blanco)
    c_dark_charcoal = RGBColor(28, 33, 41)  # Gris oscuro principal
    c_light_warm_grey = RGBColor(245, 245, 243)  # Gris claro de fondo
    c_white = RGBColor(255, 255, 255)
    c_orange = RGBColor(234, 88, 12)  # Naranjo brillante
    c_terracotta = RGBColor(154, 52, 18)  # Terracota
    c_text_dark = RGBColor(15, 23, 42)
    c_text_muted = RGBColor(100, 116, 139)
    c_border = RGBColor(226, 232, 240)
    
    image_dir = "/Users/ricardomarimodinger/.gemini/antigravity/scratch/ricardo-ai-system/images"
    
    # Helper: Set slide background color
    def set_bg_color(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Create clean card layout
    def add_card(slide, x, y, w, h, fill_color=c_white, border_color=c_border, accent_color=None):
        # Base Card
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
            
        # Top Accent Line
        if accent_color:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.12))
            accent.fill.solid()
            accent.fill.fore_color.rgb = accent_color
            accent.line.fill.background()
            
        return card

    # Helper: Add Slide Header
    def add_header(slide, title, category="PROYECTO SEGURIDAD VECINAL", is_dark=False):
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category/Badge
        p_cat = tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(13)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_orange if is_dark else c_terracotta
        p_cat.space_after = Pt(4)
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(38)
        p_title.font.bold = True
        p_title.font.color.rgb = c_white if is_dark else c_text_dark
        
        # Accent Line under header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(11.83), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = c_terracotta if not is_dark else c_orange
        line.line.fill.background()

    # Helper: Add speaker notes
    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    # Helper: Insert image safely with border card
    def insert_image_with_card(slide, x, y, w, h, img_name, accent_color=c_orange):
        add_card(slide, x, y, w, h, accent_color=accent_color)
        img_path = os.path.join(image_dir, img_name)
        if os.path.exists(img_path):
            # Inset image slightly for clean card borders
            inset = 0.15
            slide.shapes.add_picture(
                img_path, 
                Inches(x + inset), 
                Inches(y + inset + 0.1),  # Add offset for top accent
                Inches(w - (inset * 2)), 
                Inches(h - (inset * 2) - 0.1)
            )
        else:
            # Fallback text box if image missing
            tb = slide.shapes.add_textbox(Inches(x), Inches(y + 2), Inches(w), Inches(1))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"[Foto: {img_name}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.color.rgb = c_text_muted

    # Slide Layout: Blank
    blank_layout = prs.slide_layouts[6]

    # ====================================================
    # SLIDE 1: Portada (Dark charcoal background)
    s1 = prs.slides.add_slide(blank_layout)
    set_bg_color(s1, c_dark_charcoal)
    
    # Left Terracotta Side Block Accent
    block = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    block.fill.solid()
    block.fill.fore_color.rgb = c_terracotta
    block.line.fill.background()

    # Main text box
    tb_s1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf_s1 = tb_s1.text_frame
    tf_s1.word_wrap = True
    tf_s1.margin_left = tf_s1.margin_top = tf_s1.margin_right = tf_s1.margin_bottom = 0
    
    p_badge = tf_s1.paragraphs[0]
    p_badge.text = "FONDO SEGURIDAD CIUDADANA  •  CÓDIGO 2024SC0126"
    p_badge.font.name = "Arial"
    p_badge.font.size = Pt(15)
    p_badge.font.bold = True
    p_badge.font.color.rgb = c_orange
    p_badge.space_after = Pt(12)
    
    p_title = tf_s1.add_paragraph()
    p_title.text = "CUIDANDO NUESTRA SEDE"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(58)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.space_after = Pt(6)
    
    p_sub = tf_s1.add_paragraph()
    p_sub.text = "Club de Adulto Mayor Nuevo Amanecer"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = c_text_muted
    p_sub.space_after = Pt(45)
    
    p_comp = tf_s1.add_paragraph()
    p_comp.text = "ROUTER  •  Seguridad y Conectividad"
    p_comp.font.name = "Arial"
    p_comp.font.size = Pt(19)
    p_comp.font.bold = True
    p_comp.font.color.rgb = c_white
    p_comp.space_after = Pt(6)
    
    p_names = tf_s1.add_paragraph()
    p_names.text = "Presentan: Ricardo Mari  &  Ivan Santos"
    p_names.font.name = "Arial"
    p_names.font.size = Pt(16)
    p_names.font.color.rgb = c_text_muted
    
    add_notes(s1, "Muy buenas tardes a todos los socios y socias de nuestro Club de Adulto Mayor Nuevo Amanecer, a nuestra directiva y a los representantes del GORE Biobío. Es un orgullo presentarles de forma oficial y detallada el cierre del proyecto 'Cuidando Nuestra Sede', código 2024SC0126, financiado por el GORE Biobío y desarrollado por nuestra empresa Router.")

    # ====================================================
    # SLIDE 2: El Camino Hacia la Seguridad (Light Background)
    s2 = prs.slides.add_slide(blank_layout)
    set_bg_color(s2, c_light_warm_grey)
    add_header(s2, "El Sueño de una Sede Más Segura", category="Origen del Proyecto")
    
    # Left Column: White Card
    add_card(s2, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s2_l = s2.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s2_l = tb_s2_l.text_frame
    tf_s2_l.word_wrap = True
    
    p_s2_l1 = tf_s2_l.paragraphs[0]
    p_s2_l1.text = "La Necesidad Comunitaria"
    p_s2_l1.font.size = Pt(26)
    p_s2_l1.font.bold = True
    p_s2_l1.font.color.rgb = c_terracotta
    p_s2_l1.space_after = Pt(14)
    
    p_s2_l2 = tf_s2_l.add_paragraph()
    p_s2_l2.text = "• Proteger la sede social, punto de encuentro clave para nuestros adultos mayores.\n" \
                  "• Implementar disuasión efectiva frente a robos e ingresos no autorizados.\n" \
                  "• Consolidar la organización del club a través de la capacitación y tecnología de seguridad.\n" \
                  "• Proveer un entorno seguro para talleres, reuniones y actividades recreativas."
    p_s2_l2.font.size = Pt(19)
    p_s2_l2.font.color.rgb = c_text_dark
    p_s2_l2.space_after = Pt(10)
    
    # Right Column: Visual Photo Card
    insert_image_with_card(s2, 6.8, 1.8, 5.7, 4.8, "real_playground.jpg", accent_color=c_orange)
    
    add_notes(s2, "Todo gran proyecto nace de una necesidad sentida por la comunidad. En nuestro Club de Adulto Mayor, la sede es nuestro segundo hogar. El anhelo de reunirnos con tranquilidad, de realizar talleres sin miedo a robos en la sede y de proteger nuestro equipamiento, fue el motor principal. Nos dimos cuenta de que podíamos marcar una gran diferencia utilizando la tecnología a nuestro favor de forma organizada.")

    # ====================================================
    # SLIDE 3: La Postulación y la Larga Espera
    s3 = prs.slides.add_slide(blank_layout)
    set_bg_color(s3, c_light_warm_grey)
    add_header(s3, "Crónica de un Logro: Postulación y Espera", category="Cronología de Adjudicación")
    
    # 3 Column Timeline Cards
    add_card(s3, 0.75, 1.8, 3.7, 4.8, accent_color=c_text_muted)
    add_card(s3, 4.81, 1.8, 3.7, 4.8, accent_color=c_terracotta)
    add_card(s3, 8.87, 1.8, 3.7, 4.8, accent_color=c_orange)
    
    # Content Card 1
    tb_s3_c1 = s3.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(3.3), Inches(4.2))
    tf_c1 = tb_s3_c1.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "1. Postulación (2024)"
    p_c1_t.font.size = Pt(24)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = c_text_dark
    p_c1_t.space_after = Pt(10)
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = "En Abril de 2024, el Club de Adulto Mayor Nuevo Amanecer formaliza la postulación al fondo de Seguridad Ciudadana FNDR 8%.\n\nComienza el largo proceso administrativo de revisión por el Gobierno Regional."
    p_c1_b.font.size = Pt(18)
    p_c1_b.font.color.rgb = c_text_dark

    # Content Card 2
    tb_s3_c2 = s3.shapes.add_textbox(Inches(5.01), Inches(2.1), Inches(3.3), Inches(4.2))
    tf_c2 = tb_s3_c2.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "2. La Larga Espera"
    p_c2_t.font.size = Pt(24)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = c_terracotta
    p_c2_t.space_after = Pt(10)
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = "Retrasos gubernamentales extienden la espera hasta fines de 2025.\n\nRouter actúa como soporte técnico constante, aclarando plazos y dando tranquilidad a la directiva."
    p_c2_b.font.size = Pt(18)
    p_c2_b.font.color.rgb = c_text_dark

    # Content Card 3
    tb_s3_c3 = s3.shapes.add_textbox(Inches(9.07), Inches(2.1), Inches(3.3), Inches(4.2))
    tf_c3 = tb_s3_c3.text_frame
    tf_c3.word_wrap = True
    p_c3_t = tf_c3.paragraphs[0]
    p_c3_t.text = "3. Adjudicación y Éxito"
    p_c3_t.font.size = Pt(24)
    p_c3_t.font.bold = True
    p_c3_t.font.color.rgb = c_orange
    p_c3_t.space_after = Pt(10)
    p_c3_b = tf_c3.add_paragraph()
    p_c3_b.text = "La confianza en Router permite aportar documentación ágilmente ante nuevos requisitos exigidos por el GORE a última hora.\n\nFondos aprobados a fines de 2025 y transferidos en Enero de 2026."
    p_c3_b.font.size = Pt(18)
    p_c3_b.font.color.rgb = c_text_dark
    
    add_notes(s3, "Postulamos en abril de 2024. La espera administrativa del Gobierno Regional fue bastante larga y generó inquietud. En todo momento, el equipo técnico y de ingeniería de Router estuvo al lado de la directiva y los socios, explicando los plazos administrativos y manteniendo la calma. Esto cimentó una tremenda confianza que nos permitió presentar documentación ágilmente cuando el GORE exigió nuevos requisitos de última hora. Los fondos se adjudicaron a fines de 2025 y se transfirieron en Enero de 2026.")

    # ====================================================
    # SLIDE 4: Composición y Distribución del Fondo
    s4 = prs.slides.add_slide(blank_layout)
    set_bg_color(s4, c_light_warm_grey)
    add_header(s4, "¿Cómo se Invirtieron los Recursos del Fondo?", category="Presupuesto FNDR")
    
    # Left Column: White Card with text breakdown
    add_card(s4, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s4 = s4.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s4 = tb_s4.text_frame
    tf_s4.word_wrap = True
    
    p_s4_1 = tf_s4.paragraphs[0]
    p_s4_1.text = "Distribución del Presupuesto Total ($3.900.000)"
    p_s4_1.font.size = Pt(26)
    p_s4_1.font.bold = True
    p_s4_1.font.color.rgb = c_terracotta
    p_s4_1.space_after = Pt(14)
    
    p_s4_2 = tf_s4.add_paragraph()
    p_s4_2.text = "• EQUIPAMIENTO (82.1%): $3.200.000 CLP\n" \
                 "  Adquisición e instalación de las 8 cámaras IP, grabador NVR, disco duro y cableado.\n\n" \
                 "• GASTOS GENERALES (10.3%): $400.000 CLP\n" \
                 "  Coffee breaks/colaciones ($150.000), artículos de librería ($100.000) y letrero de difusión obligatorio ($150.000).\n\n" \
                 "• RECURSOS HUMANOS (7.6%): $300.000 CLP\n" \
                 "  Honorarios del profesional Iván Santos a cargo de la capacitación comunitaria."
    p_s4_2.font.size = Pt(18)
    p_s4_2.font.color.rgb = c_text_dark
    
    # Right Column: Native editable Pie/Doughnut Chart
    add_card(s4, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    chart_data = CategoryChartData()
    chart_data.categories = ['Equipamiento (82.1%)', 'Gastos Generales (10.3%)', 'Recursos Humanos (7.6%)']
    chart_data.add_series('Inversión CLP', (3200000, 400000, 300000))
    
    x, y, cx, cy = Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.0)
    chart = s4.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(14)
    chart.legend.font.name = "Arial"
    
    add_notes(s4, "Este proyecto fue financiado con un monto total ejecutado de $3.900.000 pesos chilenos. Siguiendo las directrices del GORE Biobío, el presupuesto se dividió de manera óptima para maximizar el resguardo técnico de nuestra sede social. El 82.1%, correspondiente a $3.200.000, se destinó íntegramente a la adquisición de la tecnología de vigilancia e instalación física. El 10.3% financió los gastos generales del proyecto, que incluyen coffee breaks para nuestras charlas, materiales de librería y el letrero de difusión obligatorio. Finalmente, el 7.6% restante financió la capacitación técnica a cargo del experto Iván Santos.")

    # ====================================================
    # SLIDE 5: Capacitación y Prevención (Talleres)
    s5 = prs.slides.add_slide(blank_layout)
    set_bg_color(s5, c_light_warm_grey)
    add_header(s5, "El Factor Humano: Capacitación y Prevención", category="Factor Humano y Social")
    
    # Left Card
    add_card(s5, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s5_l = s5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s5_l = tb_s5_l.text_frame
    tf_s5_l.word_wrap = True
    
    p_s5_l1 = tf_s5_l.paragraphs[0]
    p_s5_l1.text = "Charlas de Seguridad y Prevención"
    p_s5_l1.font.size = Pt(26)
    p_s5_l1.font.bold = True
    p_s5_l1.font.color.rgb = c_terracotta
    p_s5_l1.space_after = Pt(14)
    
    p_s5_l2 = tf_s5_l.add_paragraph()
    p_s5_l2.text = "• Experto Técnico a Cargo:\n" \
                  "  Ingeniero de Ejecución en Informática e instalador Iván Santos ($300.000 CLP en total).\n\n" \
                  "• Capacitación Comunitaria:\n" \
                  "  Charlas formativas enfocadas en autocuidado, prevención del delito en el entorno de la sede y uso responsable del sistema de grabaciones.\n\n" \
                  "• Cohesión y Participación:\n" \
                  "  Asistencia masiva de los socios, acompañada de coffee breaks y materiales de librería para las dinámicas grupales."
    p_s5_l2.font.size = Pt(18)
    p_s5_l2.font.color.rgb = c_text_dark
    
    # Right Card
    insert_image_with_card(s5, 6.8, 1.8, 5.7, 4.8, "vecinos_alessandri.png", accent_color=c_orange)
    
    add_notes(s5, "Para asegurar que la tecnología sea realmente útil, se contempló la capacitación a los socios. Estas actividades estuvieron a cargo de don Iván Santos, ingeniero en informática y experto en seguridad ($300.000 CLP en total). Su rol fue educar a los adultos mayores en el autocuidado, prevención y uso legal y responsable de las grabaciones del sistema, acompañando las charlas con coffee breaks y materiales de librería financiados por el fondo.")

    # ====================================================
    # SLIDE 6: Equipamiento Tecnológico Adquirido
    s6 = prs.slides.add_slide(blank_layout)
    set_bg_color(s6, c_light_warm_grey)
    add_header(s6, "Equipamiento Tecnológico de Alta Gama", category="Ficha Técnica de Equipos")
    
    # Left: Specs Grid Card
    add_card(s6, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s6_l = s6.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s6_l = tb_s6_l.text_frame
    tf_s6_l.word_wrap = True
    
    p_s6_l1 = tf_s6_l.paragraphs[0]
    p_s6_l1.text = "Ficha Técnica: Cámaras IP Bullet"
    p_s6_l1.font.size = Pt(26)
    p_s6_l1.font.bold = True
    p_s6_l1.font.color.rgb = c_terracotta
    p_s6_l1.space_after = Pt(14)
    
    p_s6_l2 = tf_s6_l.add_paragraph()
    p_s6_l2.text = "• Alta Resolución: Imágenes nítidas de 4 Megapíxeles.\n" \
                  "• Visión Nocturna Inteligente: Infrarrojo Smart Hybrid Light.\n" \
                  "• Micrófono Incorporado: Registro de audio ambiental en puntos de acceso.\n" \
                  "• Certificación Exterior IP67: Chasis metálico antivandálico resistente a lluvia y humedad.\n" \
                  "• Almacenamiento Centralizado: NVR profesional y disco duro de televigilancia de 2TB continuo.\n" \
                  "• Enlace Robusto: Cableado estructurado blindado CAT6 y Switch POE."
    p_s6_l2.font.size = Pt(19)
    p_s6_l2.font.color.rgb = c_text_dark
    
    # Right Card: Real camera photo
    insert_image_with_card(s6, 6.8, 1.8, 5.7, 4.8, "real_camera_red_wall.jpg", accent_color=c_orange)
    
    add_notes(s6, "Seleccionamos cámaras IP Bullet de calidad profesional con una resolución de 4 Megapíxeles, que brindan imágenes de alta definición. Cuentan con visión nocturna inteligente Smart Hybrid Light, micrófonos integrados para audio de seguridad, y chasis metálico con certificación IP67 contra la intemperie. Todo el cableado es blindado de cobre puro Categoría 6, conectado a un grabador NVR con disco duro de 2TB especial para televigilancia continua, evitando pérdidas de rendimiento.")

    # ====================================================
    # SLIDE 7: Diseño Local Sin Streaming
    s7 = prs.slides.add_slide(blank_layout)
    set_bg_color(s7, c_light_warm_grey)
    add_header(s7, "Diseño Técnico Local: Cero Streaming", category="Infraestructura de Red")
    
    # Left Big Card
    add_card(s7, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s7_l = s7.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s7_l = tb_s7_l.text_frame
    tf_s7_l.word_wrap = True
    p_s7_l1 = tf_s7_l.paragraphs[0]
    p_s7_l1.text = "Red e Infraestructura Local"
    p_s7_l1.font.size = Pt(26)
    p_s7_l1.font.bold = True
    p_s7_l1.font.color.rgb = c_terracotta
    p_s7_l1.space_after = Pt(12)
    p_s7_l2 = tf_s7_l.add_paragraph()
    p_s7_l2.text = "• 8 Cámaras Internas y Perimetrales: Instalación física 100% concentrada en el edificio de la sede social.\n\n" \
                  "• Cero Radioenlaces / Streaming: Al estar todos los equipos dentro del mismo inmueble, no se requirió tecnología streaming de transmisión inalámbrica a grandes distancias.\n\n" \
                  "• Conectividad Directa y Segura: Enlaces digitales físicos vía cableado estructurado CAT6 al NVR central.\n\n" \
                  "• Mayor Robustez: Sistema inmune a interferencias climáticas externas y sin costos de enlaces ni suscripciones."
    p_s7_l2.font.size = Pt(18)
    p_s7_l2.font.color.rgb = c_text_dark
    
    # Right Column: Visual Box
    insert_image_with_card(s7, 6.8, 1.8, 5.7, 4.8, "red_streaming.png", accent_color=c_orange)
 
    add_notes(s7, "Una característica muy importante del proyecto es que la instalación es local. Al concentrar las 8 cámaras dentro y en el perímetro inmediato de la sede social, no tuvimos la necesidad de utilizar tecnología de streaming o radioenlaces inalámbricos a gran distancia como en otras poblaciones. Todo se conecta por cable físico de red de alta velocidad al NVR local, lo que hace al sistema mucho más robusto, inmune a las lluvias de Lebu y libre de cualquier cobro por enlaces de datos externos.")

    # ====================================================
    # SLIDE 8: Distribución de Cámaras en la Sede
    s_map = prs.slides.add_slide(blank_layout)
    set_bg_color(s_map, c_light_warm_grey)
    add_header(s_map, "Distribución de Cámaras: Cobertura Interna", category="Mapa de Cobertura")
    
    # Left Card
    add_card(s_map, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s_map_l = s_map.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s_map_l = tb_s_map_l.text_frame
    tf_s_map_l.word_wrap = True
    
    p_s_map_l1 = tf_s_map_l.paragraphs[0]
    p_s_map_l1.text = "Resguardo de la Sede Social"
    p_s_map_l1.font.size = Pt(26)
    p_s_map_l1.font.bold = True
    p_s_map_l1.font.color.rgb = c_terracotta
    p_s_map_l1.space_after = Pt(14)
    
    p_s_map_l2 = tf_s_map_l.add_paragraph()
    p_s_map_l2.text = "• 8 Cámaras de Seguridad IP:\n" \
                      "  Distribuidas estratégicamente para cubrir los accesos principales de la sede social y perímetros vulnerables.\n\n" \
                      "• Sectores Protegidos:\n" \
                      "  Monitoreo del salón de reuniones, accesos frontales, laterales, cocina, áreas comunes y patio interior.\n\n" \
                      "• Cero Puntos Ciegos en Sede:\n" \
                      "  Permite a la directiva y socios tener un control integral del recinto para el cuidado de los insumos y la infraestructura."
    p_s_map_l2.font.size = Pt(18)
    p_s_map_l2.font.color.rgb = c_text_dark
    
    # Right Column: Visual Box
    insert_image_with_card(s_map, 6.8, 1.8, 5.7, 4.8, "mapa_camaras.jpg", accent_color=c_orange)

    add_notes(s_map, "Este es el mapa oficial de distribución de nuestras 8 cámaras de seguridad en la sede social. Como pueden ver, se cubren de manera exhaustiva todos los puntos de interés: accesos, salones, áreas comunes, cocina y patio de luz. De esta forma, evitamos que queden zonas vulnerables y protegemos de manera eficiente el espacio donde se guardan los insumos de los talleres del club.")

    # ====================================================
    # SLIDE 9: Visualización en Vivo
    s_live = prs.slides.add_slide(blank_layout)
    set_bg_color(s_live, c_light_warm_grey)
    add_header(s_live, "Monitoreo en Tiempo Real: Áreas de la Sede", category="Visualización en Vivo")
    
    # Left Card
    add_card(s_live, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s_live_l = s_live.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s_live_l = tb_s_live_l.text_frame
    tf_s_live_l.word_wrap = True
    
    p_s_live_l1 = tf_s_live_l.paragraphs[0]
    p_s_live_l1.text = "Infraestructura Local Monitoreada"
    p_s_live_l1.font.size = Pt(26)
    p_s_live_l1.font.bold = True
    p_s_live_l1.font.color.rgb = c_terracotta
    p_s_live_l1.space_after = Pt(14)
    
    p_s_live_l2 = tf_s_live_l.add_paragraph()
    p_s_live_l2.text = "• Visualización en Sede:\n" \
                       "  Pantalla local instalada de manera segura para monitoreo del recinto.\n\n" \
                       "• Calidad de Imagen HD:\n" \
                       "  Imágenes diurnas y nocturnas con gran nitidez en pasillos y accesos.\n\n" \
                       "• Visualización Móvil P2P:\n" \
                       "  Acceso configurado en los teléfonos inteligentes de la directiva autorizada."
    p_s_live_l2.font.size = Pt(19)
    p_s_live_l2.font.color.rgb = c_text_dark
    
    # Right Column: Visual Box
    insert_image_with_card(s_live, 6.8, 1.8, 5.7, 4.8, "vista_parque.jpg", accent_color=c_orange)

    add_notes(s_live, "Para el monitoreo del sistema, instalamos una pantalla local de forma segura en la sede, y configuramos la aplicación en los teléfonos de la directiva y socios autorizados por la directiva. Esto les permite revisar las cámaras en vivo y grabaciones históricas de forma remota y segura, con tecnología P2P sin costos de servicio.")

    # ====================================================
    # SLIDE 10: Claridad Tecnológica: ¿Por qué No Leen Patentes?
    s8 = prs.slides.add_slide(blank_layout)
    set_bg_color(s8, c_light_warm_grey)
    add_header(s8, "Transparencia: Límites Técnicos del Equipamiento", category="Límites Técnicos del Proyecto")
    
    # 2 Column Cost comparison layout
    add_card(s8, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    add_card(s8, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    
    # Column 1: Our cameras
    tb_s8_c1 = s8.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s8_c1 = tb_s8_c1.text_frame
    tf_s8_c1.word_wrap = True
    tf_s8_c1.paragraphs[0].text = "Cámaras 4 Megapíxeles Adquiridas"
    tf_s8_c1.paragraphs[0].font.bold = True
    tf_s8_c1.paragraphs[0].font.size = Pt(24)
    tf_s8_c1.paragraphs[0].font.color.rgb = c_terracotta
    tf_s8_c1.paragraphs[0].space_after = Pt(14)
    p_s8_c1_b = tf_s8_c1.add_paragraph()
    p_s8_c1_b.text = "• Cobertura Local: Resguardo de 8 zonas estratégicas dentro y fuera de la sede social.\n\n" \
                    "• Capacidad Técnica: Captura de rostros, vestimentas, siluetas e incidentes en general con alto detalle.\n\n" \
                    "• Criterio de Selección: Priorizar la seguridad completa del recinto del Club de Adulto Mayor."
    p_s8_c1_b.font.size = Pt(19)
    p_s8_c1_b.font.color.rgb = c_text_dark

    # Column 2: LPR cameras
    tb_s8_c2 = s8.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s8_c2 = tb_s8_c2.text_frame
    tf_s8_c2.word_wrap = True
    tf_s8_c2.paragraphs[0].text = "Lector de Patentes (Tecnología LPR)"
    tf_s8_c2.paragraphs[0].font.bold = True
    tf_s8_c2.paragraphs[0].font.size = Pt(24)
    tf_s8_c2.paragraphs[0].font.color.rgb = c_orange
    tf_s8_c2.paragraphs[0].space_after = Pt(14)
    p_s8_c2_b = tf_s8_c2.add_paragraph()
    p_s8_c2_b.text = "• No Incluido: Tecnología LPR ausente en este proyecto.\n\n" \
                    "• Factor de Costo: Un solo sistema lector de patentes LPR real y homologado cuesta 10 veces más que una cámara estándar.\n\n" \
                    "• Impacto del Lector: Elegir LPR habría agotado el presupuesto completo en vigilar una sola zona exterior, dejando la sede social desprotegida."
    p_s8_c2_b.font.size = Pt(19)
    p_s8_c2_b.font.color.rgb = c_text_dark
    
    add_notes(s8, "Queremos hacer una aclaración técnica para ser 100% transparentes: estas cámaras NO tienen lector de patentes automático. Un lector de patentes homologado (tecnología LPR) cuesta hasta 10 veces más que una cámara IP Bullet estándar. Si hubiésemos elegido LPR, el presupuesto solo nos habría alcanzado para vigilar el portón de la sede social, dejando el salón de reuniones, la cocina y los perímetros desprotegidos. Al elegir cámaras estándar de 4MP de alta definición, pudimos colocar 8 cámaras y proteger la sede completa.")

    # ====================================================
    # SLIDE 11: Grabación y Almacenamiento Optimizado
    s9 = prs.slides.add_slide(blank_layout)
    set_bg_color(s9, c_light_warm_grey)
    add_header(s9, "Estrategia de Grabación para Máxima Utilidad", category="Estrategia de Almacenamiento")
    
    # Two Columns: Left text, Right Graphic Box
    add_card(s9, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s9_l = s9.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s9_l = tb_s9_l.text_frame
    tf_s9_l.word_wrap = True
    
    p_s9_l1 = tf_s9_l.paragraphs[0]
    p_s9_l1.text = "Calidad sobre Cantidad de Días"
    p_s9_l1.font.size = Pt(26)
    p_s9_l1.font.bold = True
    p_s9_l1.font.color.rgb = c_terracotta
    p_s9_l1.space_after = Pt(14)
    
    p_s9_l2 = tf_s9_l.add_paragraph()
    p_s9_l2.text = "• Capacidad Física: Disco duro local de 2TB de alta resistencia, especial para televigilancia continua 24/7.\n\n" \
                  "• Rango Configurado: Grabación continua de los últimos 5 días históricos.\n\n" \
                  "• Razón de Ingeniería: Al concentrar los 2TB en 5 días (en lugar de comprimir a 15 días en baja calidad), las cámaras graban a máxima definición y FPS. Esto asegura evidencias nítidas y utilizables legalmente."
    p_s9_l2.font.size = Pt(19)
    p_s9_l2.font.color.rgb = c_text_dark

    # Right Card: Native Chart showing 15 Days vs 5 Days resolution quality index
    add_card(s9, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    chart_data_s9 = CategoryChartData()
    chart_data_s9.categories = ['Configurado 5 días', 'Configurado 15 días']
    chart_data_s9.add_series('Índice de Nitidez (1-100)', (95, 30))
    chart_data_s9.add_series('FPS (Fotogramas x seg)', (25, 8))
    
    x_s9, y_s9, cx_s9, cy_s9 = Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.0)
    chart_s9 = s9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x_s9, y_s9, cx_s9, cy_s9, chart_data_s9
    ).chart
    chart_s9.has_legend = True
    chart_s9.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart_s9.legend.font.size = Pt(14)
    
    add_notes(s9, "¿Cómo funciona el almacenamiento de las grabaciones? El grabador NVR cuenta con un disco duro de 2 Terabytes de alta resistencia. Técnicamente, las cámaras podrían grabar hasta 15 días continuos en baja calidad comprimida. Decidimos configurarlas para almacenar exactamente 5 días de historial. Esto nos permite grabar en la máxima definición de los lentes sin ralentizaciones, garantizando que el video sirva como evidencia nítida e indiscutible si ocurre algún hecho en la sede social.")

    # ====================================================
    # SLIDE 12: Responsabilidad Local: Encargado de Cámaras
    s10 = prs.slides.add_slide(blank_layout)
    set_bg_color(s10, c_light_warm_grey)
    add_header(s10, "Custodia de la Sede y Resguardo del Sistema", category="Administración del Sistema")
    
    # Left Column: Info Card
    add_card(s10, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s10_l = s10.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s10_l = tb_s10_l.text_frame
    tf_s10_l.word_wrap = True
    
    p_s10_l1 = tf_s10_l.paragraphs[0]
    p_s10_l1.text = "Custodia y Nexo de Seguridad"
    p_s10_l1.font.size = Pt(26)
    p_s10_l1.font.bold = True
    p_s10_l1.font.color.rgb = c_terracotta
    p_s10_l1.space_after = Pt(14)
    
    p_s10_l2 = tf_s10_l.add_paragraph()
    p_s10_l2.text = "• Custodia de la Directiva: El NVR y la pantalla de monitoreo están resguardados bajo llave en la sede social.\n\n" \
                   "• Control de Acceso: Solo los directivos del Club de Adulto Mayor tienen acceso directo a la central física de los equipos.\n\n" \
                   "• Privacidad de Datos: Extracción de clips confidencial y restringida, respetando a los socios y resguardando la sede."
    p_s10_l2.font.size = Pt(19)
    p_s10_l2.font.color.rgb = c_text_dark
    
    # Right Column: Photo Card
    insert_image_with_card(s10, 6.8, 1.8, 5.7, 4.8, "real_neighbors.jpg", accent_color=c_orange)

    add_notes(s10, "Un sistema de seguridad en una sede comunitaria requiere una administración muy cuidadosa de los equipos. El grabador y la pantalla se resguardan en un gabinete bajo llave dentro de la sede social. La directiva del club, encabezada por nuestra Presidenta Isolda Camaño, ejerce la custodia física de los equipos, garantizando que nadie manipule la central y que se respete de forma estricta la privacidad de todas las personas.")

    # ====================================================
    # SLIDE 13: Protocolo de Seguridad y Privacidad
    s11 = prs.slides.add_slide(blank_layout)
    set_bg_color(s11, c_light_warm_grey)
    add_header(s11, "Protocolo de Seguridad y Privacidad", category="Protocolo de Privacidad")
    
    # Left Card
    add_card(s11, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s11_l = s11.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s11_l = tb_s11_l.text_frame
    tf_s11_l.word_wrap = True
    p_s11_l1 = tf_s11_l.paragraphs[0]
    p_s11_l1.text = "¿Cuándo se Activa el Protocolo?"
    p_s11_l1.font.size = Pt(26)
    p_s11_l1.font.bold = True
    p_s11_l1.font.color.rgb = c_terracotta
    p_s11_l1.space_after = Pt(14)
    p_s11_l2 = tf_s11_l.add_paragraph()
    p_s11_l2.text = "• Eventos en la Sede: Sospecha o comisión de delitos (robos, hurtos o ingresos de terceros a la sede).\n\n" \
                   "• Daños e Incivilidades: Vandalismo sobre fachadas, rejas o infraestructura del Club.\n\n" \
                   "• Evidencias Legales: Apoyo a investigaciones policiales oficiales en la zona inmediata."
    p_s11_l2.font.size = Pt(19)
    p_s11_l2.font.color.rgb = c_text_dark

    # Right Card
    add_card(s11, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    tb_s11_r = s11.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s11_r = tb_s11_r.text_frame
    tf_s11_r.word_wrap = True
    p_s11_r1 = tf_s11_r.paragraphs[0]
    p_s11_r1.text = "Límites Técnicos y de Privacidad"
    p_s11_r1.font.size = Pt(26)
    p_s11_r1.font.bold = True
    p_s11_r1.font.color.rgb = c_orange
    p_s11_r1.space_after = Pt(14)
    p_s11_r2 = tf_s11_r.add_paragraph()
    p_s11_r2.text = "• Enfoque de Cámaras: Calibración específica para registrar accesos y salones multiuso del Club.\n\n" \
                   "• Resguardo Colindante: Se evitan campos de visión hacia ventanas o patios privados de las viviendas vecinas.\n\n" \
                   "• Uso del Sistema: Las grabaciones tienen uso exclusivamente de seguridad vecinal, no con fines recreativos."
    p_s11_r2.font.size = Pt(19)
    p_s11_r2.font.color.rgb = c_text_dark

    add_notes(s11, "Para resguardar el derecho a la privacidad de todos los hogares vecinos y de nuestros socios, diseñamos un protocolo transparente de acceso a las grabaciones. Las cámaras están hechas para proteger. Si un vecino o socio necesita revisar una grabación debido a un robo o daño, debe hacer una solicitud formal a la Directiva de la Junta. La directiva analizará el caso de forma confidencial y coordinará la extracción del clip para que sea enviado vía WhatsApp al afectado o entregado directamente a las policías como medio de prueba legal.")

    # ====================================================
    # SLIDE 14: Pasos Formales de Solicitud (Protocolo)
    s12 = prs.slides.add_slide(blank_layout)
    set_bg_color(s12, c_light_warm_grey)
    add_header(s12, "Pasos Formales de Solicitud (Protocolo)", category="Protocolo Formal")
    
    # 4 Grid Steps
    add_card(s12, 0.75, 1.8, 5.7, 2.3, accent_color=c_terracotta)
    add_card(s12, 6.8, 1.8, 5.7, 2.3, accent_color=c_orange)
    add_card(s12, 0.75, 4.3, 5.7, 2.3, accent_color=c_text_muted)
    add_card(s12, 6.8, 4.3, 5.7, 2.3, accent_color=c_dark_charcoal)
    
    # Step 1
    tb_s12_s1 = s12.shapes.add_textbox(Inches(0.95), Inches(2.05), Inches(5.3), Inches(1.9))
    tf_s12_s1 = tb_s12_s1.text_frame
    tf_s12_s1.word_wrap = True
    tf_s12_s1.paragraphs[0].text = "Paso 1: Solicitud formal a la Directiva"
    tf_s12_s1.paragraphs[0].font.bold = True
    tf_s12_s1.paragraphs[0].font.size = Pt(20)
    tf_s12_s1.paragraphs[0].font.color.rgb = c_terracotta
    p_s12_s1_b = tf_s12_s1.add_paragraph()
    p_s12_s1_b.text = "El socio formaliza su solicitud a la directiva por escrito o a través del WhatsApp de la Presidenta. No debe acudir directamente al cuidador técnico."
    p_s12_s1_b.font.size = Pt(16)
    p_s12_s1_b.font.color.rgb = c_text_dark
    p_s12_s1_b.space_before = Pt(4)

    # Step 2
    tb_s12_s2 = s12.shapes.add_textbox(Inches(7.0), Inches(2.05), Inches(5.3), Inches(1.9))
    tf_s12_s2 = tb_s12_s2.text_frame
    tf_s12_s2.word_wrap = True
    tf_s12_s2.paragraphs[0].text = "Paso 2: Datos Técnicos del Incidente"
    tf_s12_s2.paragraphs[0].font.bold = True
    tf_s12_s2.paragraphs[0].font.size = Pt(20)
    tf_s12_s2.paragraphs[0].font.color.rgb = c_orange
    p_s12_s2_b = tf_s12_s2.add_paragraph()
    p_s12_s2_b.text = "Especificar la fecha exacta del incidente, rango aproximado de hora y la zona de la sede o la cámara asociada al suceso."
    p_s12_s2_b.font.size = Pt(16)
    p_s12_s2_b.font.color.rgb = c_text_dark
    p_s12_s2_b.space_before = Pt(4)

    # Step 3
    tb_s12_s3 = s12.shapes.add_textbox(Inches(0.95), Inches(4.55), Inches(5.3), Inches(1.9))
    tf_s12_s3 = tb_s12_s3.text_frame
    tf_s12_s3.word_wrap = True
    tf_s12_s3.paragraphs[0].text = "Paso 3: Extracción y Envío de Evidencia"
    tf_s12_s3.paragraphs[0].font.bold = True
    tf_s12_s3.paragraphs[0].font.size = Pt(20)
    tf_s12_s3.paragraphs[0].font.color.rgb = c_text_muted
    p_s12_s3_b = tf_s12_s3.add_paragraph()
    p_s12_s3_b.text = "La directiva revisa y aprueba. Se extrae el fragmento de video de forma confidencial y se envía por WhatsApp o pendrive como prueba legal."
    p_s12_s3_b.font.size = Pt(16)
    p_s12_s3_b.font.color.rgb = c_text_dark
    p_s12_s3_b.space_before = Pt(4)

    # Step 4
    tb_s12_s4 = s12.shapes.add_textbox(Inches(7.0), Inches(4.55), Inches(5.3), Inches(1.9))
    tf_s12_s4 = tb_s12_s4.text_frame
    tf_s12_s4.word_wrap = True
    tf_s12_s4.paragraphs[0].text = "Límite Crítico: Plazo de 5 Días"
    tf_s12_s4.paragraphs[0].font.bold = True
    tf_s12_s4.paragraphs[0].font.size = Pt(20)
    tf_s12_s4.paragraphs[0].font.color.rgb = c_dark_charcoal
    p_s12_s4_b = tf_s12_s4.add_paragraph()
    p_s12_s4_b.text = "IMPORTANTE: Realizar el requerimiento antes de cumplirse 5 días del hecho. Posterior a esto, el almacenamiento sobrescribirá el video."
    p_s12_s4_b.font.size = Pt(16)
    p_s12_s4_b.font.color.rgb = c_text_dark
    p_s12_s4_b.space_before = Pt(4)
 
    add_notes(s12, "Un detalle muy importante que debemos memorizar: el plazo máximo para solicitar una grabación es de 5 días. Debido a que el sistema sobrescribe las grabaciones antiguas para mantener la mejor calidad y fotogramas de video, si un incidente ocurre un lunes, la directiva debe recibir la solicitud antes del viernes de esa misma semana para asegurar que el video no se pierda.")

    # ====================================================
    # SLIDE 15: Aclarando Dudas Comunitarias (FAQ)
    s13 = prs.slides.add_slide(blank_layout)
    set_bg_color(s13, c_light_warm_grey)
    add_header(s13, "Aclarando Dudas Frecuentes (FAQ)", category="Preguntas Frecuentes")
    
    # 2 columns of FAQ cards
    add_card(s13, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    add_card(s13, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    
    # FAQ Column 1
    tb_s13_c1 = s13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_s13_c1 = tb_s13_c1.text_frame
    tf_s13_c1.word_wrap = True
    
    p_faq_1 = tf_s13_c1.paragraphs[0]
    p_faq_1.text = "¿Quién financia la electricidad?"
    p_faq_1.font.bold = True
    p_faq_1.font.size = Pt(22)
    p_faq_1.font.color.rgb = c_terracotta
    p_faq_1.space_after = Pt(2)
    p_faq_1_a = tf_s13_c1.add_paragraph()
    p_faq_1_a.text = "Costo mensual aproximado de $1.200 por cámara. Este consumo eléctrico es mínimo y se asume como gasto de luz operacional básico de la sede social."
    p_faq_1_a.font.size = Pt(17)
    p_faq_1_a.font.color.rgb = c_text_dark
    p_faq_1_a.space_after = Pt(10)
    
    p_faq_2 = tf_s13_c1.add_paragraph()
    p_faq_2.text = "¿Tiene costo mensual para el socio?"
    p_faq_2.font.bold = True
    p_faq_2.font.size = Pt(22)
    p_faq_2.font.color.rgb = c_terracotta
    p_faq_2.space_after = Pt(2)
    p_faq_2_a = tf_s13_c1.add_paragraph()
    p_faq_2_a.text = "No. El fondo FNDR financió el 100% de la instalación y equipos. Se entrega en propiedad del Club de Adulto Mayor libre de cuotas mensuales fijas ($0 de mantención)."
    p_faq_2_a.font.size = Pt(17)
    p_faq_2_a.font.color.rgb = c_text_dark

    # FAQ Column 2
    tb_s13_c2 = s13.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_s13_c2 = tb_s13_c2.text_frame
    tf_s13_c2.word_wrap = True
    
    p_faq_3 = tf_s13_c2.paragraphs[0]
    p_faq_3.text = "¿Se transmiten por Internet pública?"
    p_faq_3.font.bold = True
    p_faq_3.font.size = Pt(22)
    p_faq_3.font.color.rgb = c_orange
    p_faq_3.space_after = Pt(2)
    p_faq_3_a = tf_s13_c2.add_paragraph()
    p_faq_3_a.text = "No. La red de cámaras es completamente física y local. No se suben fragmentos de video a nubes públicas ni redes de libre acceso, resguardando los datos."
    p_faq_3_a.font.size = Pt(17)
    p_faq_3_a.font.color.rgb = c_text_dark
    p_faq_3_a.space_after = Pt(10)
    
    p_faq_4 = tf_s13_c2.add_paragraph()
    p_faq_4.text = "¿Cada cuánto se hace mantención?"
    p_faq_4.font.bold = True
    p_faq_4.font.size = Pt(22)
    p_faq_4.font.color.rgb = c_orange
    p_faq_4.space_after = Pt(2)
    p_faq_4_a = tf_s13_c2.add_paragraph()
    p_faq_4_a.text = "Se recomiendan visitas de mantención preventiva cada 3 meses para limpiar lentes, comprobar anclajes y revisar el correcto funcionamiento del disco duro local."
    p_faq_4_a.font.size = Pt(17)
    p_faq_4_a.font.color.rgb = c_text_dark

    add_notes(s13, "En estas asambleas siempre surgen inquietudes lógicas. El gasto eléctrico es mínimo, y al estar conectadas al medidor de la sede social, se incorpora como parte del consumo operacional normal de luz. El sistema es propio y está pagado en su totalidad, por lo que no hay cobros mensuales ocultos ni cuotas para los socios. Los datos están resguardados físicamente en el NVR y no se suben a internet. Para garantizar la vida útil de los equipos, recomendamos una limpieza y mantención cada 3 meses.")

    # ====================================================
    # SLIDE 16: Garantía Oficial del Sistema
    s14 = prs.slides.add_slide(blank_layout)
    set_bg_color(s14, c_light_warm_grey)
    add_header(s14, "Garantía Oficial del Sistema", category="Garantía del Proyecto")
    
    # 2 columns of Warranty cards
    add_card(s14, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    add_card(s14, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    
    # Column 1: Cobertura
    tb_s14_c1 = s14.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s14_c1 = tb_s14_c1.text_frame
    tf_s14_c1.word_wrap = True
    
    p_w1 = tf_s14_c1.paragraphs[0]
    p_w1.text = "Cobertura de Garantía Router"
    p_w1.font.bold = True
    p_w1.font.size = Pt(26)
    p_w1.font.color.rgb = c_terracotta
    p_w1.space_after = Pt(14)
    
    p_w1_b = tf_s14_c1.add_paragraph()
    p_w1_b.text = "• Vigencia: 6 meses de garantía oficial a partir de la entrega conforme del proyecto (Mayo 2026).\n\n" \
                  "• Alcance Técnico: Cubre fallas técnicas del equipamiento por defectos de fabricación o fallas en el montaje e instalación de cableado."
    p_w1_b.font.size = Pt(19)
    p_w1_b.font.color.rgb = c_text_dark
    
    # Column 2: Exclusiones
    tb_s14_c2 = s14.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s14_c2 = tb_s14_c2.text_frame
    tf_s14_c2.word_wrap = True
    
    p_w2 = tf_s14_c2.paragraphs[0]
    p_w2.text = "Exclusiones de Cobertura"
    p_w2.font.bold = True
    p_w2.font.size = Pt(26)
    p_w2.font.color.rgb = c_orange
    p_w2.space_after = Pt(14)
    
    p_w2_b = tf_s14_c2.add_paragraph()
    p_w2_b.text = "• Intervención de Terceros: La garantía queda sin efecto en caso de manipulación no autorizada o modificaciones técnicas ajenas a la empresa.\n\n" \
                  "• Siniestros Excluidos: No cubre daños provocados por vandalismo, robos de equipos, sobretensiones eléctricas externas o siniestros por factores climáticos destructivos extremos."
    p_w2_b.font.size = Pt(19)
    p_w2_b.font.color.rgb = c_text_dark
    
    add_notes(s14, "Nuestra empresa Router entrega una garantía técnica oficial de 6 meses desde el acta de entrega firmada el 12 de Mayo de 2026. Cubre defectos de fábrica o de montaje. No cubre, por supuesto, la manipulación de personas externas o modificaciones no autorizadas, ni daños por robos, vandalismo, alzas de voltaje externas de la red o desastres naturales. Debemos coordinarnos como club para cuidar el sistema.")

    # ====================================================
    # SLIDE 17: Espacio de Consultas y Diálogo
    s15 = prs.slides.add_slide(blank_layout)
    set_bg_color(s15, c_light_warm_grey)
    add_header(s15, "Foro Abierto: Preguntas y Diálogo Vecinal", category="Asamblea Comunitaria")
    
    add_card(s15, 0.75, 1.8, 11.83, 4.8, accent_color=c_terracotta)
    tb_s15 = s15.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf_s15 = tb_s15.text_frame
    tf_s15.word_wrap = True
    
    p_s15_1 = tf_s15.paragraphs[0]
    p_s15_1.text = "Espacio de consultas para la asamblea"
    p_s15_1.font.size = Pt(34)
    p_s15_1.font.bold = True
    p_s15_1.font.color.rgb = c_terracotta
    p_s15_1.alignment = PP_ALIGN.CENTER
    p_s15_1.space_after = Pt(24)
    
    p_s15_2 = tf_s15.add_paragraph()
    p_s15_2.text = "Su opinión y sus dudas son muy importantes.\n\n" \
                  "Por favor, levanten la mano y compartan sus inquietudes con la Directiva del Club de Adulto Mayor y el Equipo Técnico de Router."
    p_s15_2.font.size = Pt(24)
    p_s15_2.font.color.rgb = c_text_dark
    p_s15_2.alignment = PP_ALIGN.CENTER
    
    add_notes(s15, "Ahora abrimos la palabra para todos ustedes. Queremos escuchar sus inquietudes, dudas o sugerencias respecto al sistema. Toda pregunta es bienvenida para que todos nos vayamos a casa con absoluta claridad sobre el uso y beneficio de nuestro sistema de seguridad.")

    # ====================================================
    # SLIDE 18: Difusión y Placa GORE
    s16 = prs.slides.add_slide(blank_layout)
    set_bg_color(s16, c_light_warm_grey)
    add_header(s16, "Placa Oficial FNDR y Difusión del Proyecto", category="Difusión del Proyecto")
    
    # Left Card
    add_card(s16, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s16_l = s16.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf_s16_l = tb_s16_l.text_frame
    tf_s16_l.word_wrap = True
    
    p_s16_l1 = tf_s16_l.paragraphs[0]
    p_s16_l1.text = "La Seguridad se Construye Comunicando"
    p_s16_l1.font.size = Pt(26)
    p_s16_l1.font.bold = True
    p_s16_l1.font.color.rgb = c_terracotta
    p_s16_l1.space_after = Pt(12)
    
    p_s16_l2 = tf_s16_l.add_paragraph()
    p_s16_l2.text = "• Placa Oficial Obligatoria: Instalación de una placa de acrílico exterior en la sede social que visibiliza el financiamiento del GORE Biobío ($150.000 del ítem de difusión).\n\n" \
                   "• Diseño Estandarizado: Placa con la imagen institucional y el logo oficial del GORE al costado izquierdo.\n\n" \
                   "• Foco Informativo: Folletos y afiches utilizados durante las charlas formativas para promover la corresponsabilidad vecinal."
    p_s16_l2.font.size = Pt(17)
    p_s16_l2.font.color.rgb = c_text_dark
    
    # Right Card
    insert_image_with_card(s16, 6.8, 1.8, 5.7, 4.8, "video_promo.png", accent_color=c_orange)
    
    add_notes(s16, "Una parte muy relevante del financiamiento del GORE Biobío es la difusión del proyecto. Destinamos $150.000 pesos a la fabricación e instalación de la placa acrílica exterior obligatoria y a los folletos e invitaciones de difusión. La placa exterior ya está colocada en la fachada y cuenta con el logotipo oficial del Gobierno Regional del Biobío al costado izquierdo, certificando el origen público del financiamiento.")

    # ====================================================
    # SLIDE 19: Promoción Especial Vecinos
    s_promo = prs.slides.add_slide(blank_layout)
    set_bg_color(s_promo, c_light_warm_grey)
    add_header(s_promo, "Promoción Especial para tu Hogar", category="Beneficio Comunitario")
    
    # Left Card
    add_card(s_promo, 0.75, 1.8, 5.7, 4.8, accent_color=c_orange)
    tb_sp_l = s_promo.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf_sp_l = tb_sp_l.text_frame
    tf_sp_l.word_wrap = True
    
    p_sp_l1 = tf_sp_l.paragraphs[0]
    p_sp_l1.text = "Instalación de 1 Cámara Domiciliaria"
    p_sp_l1.font.size = Pt(26)
    p_sp_l1.font.bold = True
    p_sp_l1.font.color.rgb = c_orange
    p_sp_l1.space_after = Pt(6)
    
    p_sp_l2 = tf_sp_l.add_paragraph()
    p_sp_l2.text = "Precio Preferencial: $90.000 CLP (Pago Único)"
    p_sp_l2.font.size = Pt(19)
    p_sp_l2.font.bold = True
    p_sp_l2.font.color.rgb = c_terracotta
    p_sp_l2.space_after = Pt(10)
    
    p_sp_l3 = tf_sp_l.add_paragraph()
    p_sp_l3.text = "• Cámara de Alta Gama: Full HD, visión nocturna, audio.\n" \
                   "• Almacenamiento Local: Tarjeta Micro SD de 64GB incluida.\n" \
                   "• Todo Incluido: Instalación técnica, configuración y garantía Router.\n" \
                   "• Exclusivo Socios: Promoción especial de apoyo a la seguridad familiar."
    p_sp_l3.font.size = Pt(16)
    p_sp_l3.font.color.rgb = c_text_dark
    p_sp_l3.space_after = Pt(12)
    
    # WhatsApp Highlight Box in Left Card
    promo_shape = s_promo.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.2), Inches(5.4), Inches(1.2))
    promo_shape.fill.solid()
    promo_shape.fill.fore_color.rgb = RGBColor(240, 253, 244)  # Light green tint
    promo_shape.line.color.rgb = RGBColor(22, 163, 74)
    promo_shape.line.width = Pt(1.5)
    
    tf_p = promo_shape.text_frame
    tf_p.word_wrap = True
    tf_p.margin_top = Inches(0.08)
    tf_p.margin_bottom = Inches(0.08)
    tf_p.margin_left = Inches(0.15)
    tf_p.margin_right = Inches(0.15)
    
    p1 = tf_p.paragraphs[0]
    p1.text = "¡RESERVA TU INSTALACIÓN DOMICILIARIA!"
    p1.font.bold = True
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(22, 163, 74)
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(2)
    
    p2 = tf_p.add_paragraph()
    p2.text = "WhatsApp: +56 9 7858 9090"
    p2.font.bold = True
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(22, 163, 74)
    p2.alignment = PP_ALIGN.CENTER
    
    # Right Card
    insert_image_with_card(s_promo, 6.8, 1.8, 5.7, 4.8, "camera_alessandri.png", accent_color=c_orange)
    
    add_notes(s_promo, "Como beneficio adicional para todos los socios de nuestro Club de Adulto Mayor, desde Router queremos apoyarlos con su seguridad familiar en el hogar. Ofrecemos una promoción especial para instalar una cámara domiciliaria particular Full HD, con visión nocturna, audio y tarjeta de 64GB por un pago único preferencial de $90.000 pesos, que incluye la instalación técnica y la garantía oficial de Router. Pueden agendarla al WhatsApp que aparece en pantalla.")

    # ====================================================
    # SLIDE 20: Conclusión y Agradecimientos (Dark charcoal background)
    s17 = prs.slides.add_slide(blank_layout)
    set_bg_color(s17, c_dark_charcoal)
    
    # Left Terracotta Stripe
    block_s17 = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    block_s17.fill.solid()
    block_s17.fill.fore_color.rgb = c_terracotta
    block_s17.line.fill.background()
    
    # Title & Subtitle Box
    tb_s17_top = s17.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(3.0))
    tf_s17_top = tb_s17_top.text_frame
    tf_s17_top.word_wrap = True
    tf_s17_top.margin_left = tf_s17_top.margin_top = tf_s17_top.margin_right = tf_s17_top.margin_bottom = 0
    
    p_s17_1 = tf_s17_top.paragraphs[0]
    p_s17_1.text = "¡JUNTOS CUIDAMOS NUESTRA SEDE!"
    p_s17_1.font.name = "Arial"
    p_s17_1.font.size = Pt(46)
    p_s17_1.font.bold = True
    p_s17_1.font.color.rgb = c_white
    p_s17_1.space_after = Pt(16)
    
    p_s17_2 = tf_s17_top.add_paragraph()
    p_s17_2.text = '"La tecnología vigila, pero la comunidad organizada protege."'
    p_s17_2.font.name = "Arial"
    p_s17_2.font.size = Pt(26)
    p_s17_2.font.bold = True
    p_s17_2.font.color.rgb = c_orange
    
    # Footer separator line
    line_s17 = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.2), Inches(11.0), Inches(0.02))
    line_s17.fill.solid()
    line_s17.fill.fore_color.rgb = RGBColor(71, 85, 105)
    line_s17.line.fill.background()
    
    # Left Footer Box: Agradecimientos
    tb_s17_agrad = s17.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(5.5), Inches(2.2))
    tf_s17_agrad = tb_s17_agrad.text_frame
    tf_s17_agrad.word_wrap = True
    tf_s17_agrad.margin_left = tf_s17_agrad.margin_top = tf_s17_agrad.margin_right = tf_s17_agrad.margin_bottom = 0
    
    p_agrad_title = tf_s17_agrad.paragraphs[0]
    p_agrad_title.text = "AGRADECIMIENTOS ESPECIALES:"
    p_agrad_title.font.name = "Arial"
    p_agrad_title.font.bold = True
    p_agrad_title.font.size = Pt(16)
    p_agrad_title.font.color.rgb = RGBColor(148, 163, 184)
    p_agrad_title.space_after = Pt(8)
    
    p_agrad_body = tf_s17_agrad.add_paragraph()
    p_agrad_body.text = "• Gobierno Regional del Biobío (GORE Biobío)\n" \
                        "• Directiva Club de Adulto Mayor Nuevo Amanecer\n" \
                        "• Presidenta doña Isolda Camaño Lavín"
    p_agrad_body.font.name = "Arial"
    p_agrad_body.font.size = Pt(16)
    p_agrad_body.font.color.rgb = c_text_muted
    
    # Vertical Separator between footers
    line_vert = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(4.5), Inches(0.02), Inches(1.8))
    line_vert.fill.solid()
    line_vert.fill.fore_color.rgb = RGBColor(71, 85, 105)
    line_vert.line.fill.background()
    
    # Right Footer Box: Contacto Router
    tb_s17_contact = s17.shapes.add_textbox(Inches(7.3), Inches(4.5), Inches(5.0), Inches(2.2))
    tf_s17_contact = tb_s17_contact.text_frame
    tf_s17_contact.word_wrap = True
    tf_s17_contact.margin_left = tf_s17_contact.margin_top = tf_s17_contact.margin_right = tf_s17_contact.margin_bottom = 0
    
    p_contact_title = tf_s17_contact.paragraphs[0]
    p_contact_title.text = "CONTACTO ROUTER:"
    p_contact_title.font.name = "Arial"
    p_contact_title.font.bold = True
    p_contact_title.font.size = Pt(16)
    p_contact_title.font.color.rgb = RGBColor(148, 163, 184)
    p_contact_title.space_after = Pt(8)
    
    p_contact_phone = tf_s17_contact.add_paragraph()
    p_contact_phone.text = "WhatsApp: +56 9 7858 9090"
    p_contact_phone.font.name = "Arial"
    p_contact_phone.font.bold = True
    p_contact_phone.font.size = Pt(22)
    p_contact_phone.font.color.rgb = RGBColor(22, 163, 74)
    p_contact_phone.space_after = Pt(4)
    
    p_contact_sub = tf_s17_contact.add_paragraph()
    p_contact_sub.text = "Router Ingeniería & Conectividad"
    p_contact_sub.font.name = "Arial"
    p_contact_sub.font.size = Pt(16)
    p_contact_sub.font.color.rgb = c_text_muted
    
    add_notes(s17, "Llegamos al final de esta presentación, pero esto es el inicio de una etapa mucho más segura para nuestro Club. Queremos reiterar nuestro más sincero agradecimiento al Gobierno Regional del Biobío por los fondos y su apoyo; a nuestra directiva por el trabajo constante en terreno; y a cada uno de ustedes, socios y socias, por su compromiso y participación. Como dice nuestro lema: la tecnología vigila, pero la solidaridad y organización es lo que realmente nos protege. ¡Muchas gracias a todos!")

    # Save presentation
    output_path = "/Users/ricardomarimodinger/.gemini/antigravity/scratch/ricardo-ai-system/presentacion_nuevo_amanecer.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_premium_presentation()
