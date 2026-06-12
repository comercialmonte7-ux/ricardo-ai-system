import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def create_premium_presentation():
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Palette definition (Gris, Naranjo, Terracota, Blanco)
    c_dark_charcoal = RGBColor(28, 33, 41)  # Gris oscuro principal
    c_light_warm_grey = RGBColor(245, 245, 243)  # Gris claro de fondo
    c_white = RGBColor(255, 255, 255)
    c_orange = RGBColor(234, 88, 12)  # Naranjo brillante
    c_terracotta = RGBColor(154, 52, 18)  # Terracota
    c_text_dark = RGBColor(15, 23, 42)
    c_text_muted = RGBColor(100, 116, 139)
    c_border = RGBColor(226, 232, 240)
    
    image_dir = "/Users/ricardomarimodinger/.gemini/antigravity/scratch/ricardo-ai-system/images"
    
    # Helper: Set slide background color
    def set_bg_color(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Create clean card layout
    def add_card(slide, x, y, w, h, fill_color=c_white, border_color=c_border, accent_color=None):
        # Base Card
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
            
        # Top Accent Line
        if accent_color:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.12))
            accent.fill.solid()
            accent.fill.fore_color.rgb = accent_color
            accent.line.fill.background()
            
        return card

    # Helper: Add Slide Header
    def add_header(slide, title, category="PROYECTO SEGURIDAD VECINAL", is_dark=False):
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category/Badge
        p_cat = tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(13)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_orange if is_dark else c_terracotta
        p_cat.space_after = Pt(4)
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(38)
        p_title.font.bold = True
        p_title.font.color.rgb = c_white if is_dark else c_text_dark
        
        # Accent Line under header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(11.83), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = c_terracotta if not is_dark else c_orange
        line.line.fill.background()

    # Helper: Add speaker notes
    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    # Helper: Insert image safely with border card
    def insert_image_with_card(slide, x, y, w, h, img_name, accent_color=c_orange):
        add_card(slide, x, y, w, h, accent_color=accent_color)
        img_path = os.path.join(image_dir, img_name)
        if os.path.exists(img_path):
            # Inset image slightly for clean card borders
            inset = 0.15
            slide.shapes.add_picture(
                img_path, 
                Inches(x + inset), 
                Inches(y + inset + 0.1),  # Add offset for top accent
                Inches(w - (inset * 2)), 
                Inches(h - (inset * 2) - 0.1)
            )
        else:
            # Fallback text box if image missing
            tb = slide.shapes.add_textbox(Inches(x), Inches(y + 2), Inches(w), Inches(1))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"[Foto: {img_name}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.color.rgb = c_text_muted

    # Slide Layout: Blank
    blank_layout = prs.slide_layouts[6]

    # ====================================================
    # SLIDE 1: Portada (Dark charcoal background)
    s1 = prs.slides.add_slide(blank_layout)
    set_bg_color(s1, c_dark_charcoal)
    
    # Left Terracotta Side Block Accent
    block = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    block.fill.solid()
    block.fill.fore_color.rgb = c_terracotta
    block.line.fill.background()

    # Main text box
    tb_s1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf_s1 = tb_s1.text_frame
    tf_s1.word_wrap = True
    tf_s1.margin_left = tf_s1.margin_top = tf_s1.margin_right = tf_s1.margin_bottom = 0
    
    p_badge = tf_s1.paragraphs[0]
    p_badge.text = "FONDO SEGURIDAD CIUDADANA  •  CÓDIGO 2024SCXXXX"
    p_badge.font.name = "Arial"
    p_badge.font.size = Pt(15)
    p_badge.font.bold = True
    p_badge.font.color.rgb = c_orange
    p_badge.space_after = Pt(12)
    
    p_title = tf_s1.add_paragraph()
    p_title.text = "PROTEGIENDO VISTA AL MAR"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(58)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.space_after = Pt(6)
    
    p_sub = tf_s1.add_paragraph()
    p_sub.text = "Sistema de Seguridad, Monitoreo y Prevención Comunitaria"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = c_text_muted
    p_sub.space_after = Pt(45)
    
    p_comp = tf_s1.add_paragraph()
    p_comp.text = "ROUTER  •  Tecnología y Seguridad"
    p_comp.font.name = "Arial"
    p_comp.font.size = Pt(19)
    p_comp.font.bold = True
    p_comp.font.color.rgb = c_white
    p_comp.space_after = Pt(6)
    
    p_names = tf_s1.add_paragraph()
    p_names.text = "Presentan: Ricardo Mari  &  Ivan Santos"
    p_names.font.name = "Arial"
    p_names.font.size = Pt(16)
    p_names.font.color.rgb = c_text_muted
    
    add_notes(s1, "Muy buenas tardes a todos los vecinos y vecinas, a la directiva de nuestra Junta de Vecinos y a los representantes del GORE. Es un orgullo presentarles de forma oficial y detallada el cierre del proyecto 'Protegiendo Vista al Mar', código 2024SCXXXX, financiado por el GORE Biobío y desarrollado por nuestra empresa Router.")

    # ====================================================
    # SLIDE 2: El Camino Hacia la Seguridad (Light Background)
    s2 = prs.slides.add_slide(blank_layout)
    set_bg_color(s2, c_light_warm_grey)
    add_header(s2, "El Sueño de un Barrio Más Seguro")
    
    # Left Column: White Card
    add_card(s2, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s2_l = s2.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s2_l = tb_s2_l.text_frame
    tf_s2_l.word_wrap = True
    
    p_s2_l1 = tf_s2_l.paragraphs[0]
    p_s2_l1.text = "La Necesidad Comunitaria"
    p_s2_l1.font.size = Pt(26)
    p_s2_l1.font.bold = True
    p_s2_l1.font.color.rgb = c_terracotta
    p_s2_l1.space_after = Pt(14)
    
    p_s2_l2 = tf_s2_l.add_paragraph()
    p_s2_l2.text = "• Recuperar la tranquilidad en las calles para las familias y niños.\n" \
                  "• Implementar disuasión efectiva frente a robos e incivilidades.\n" \
                  "• Consolidar la organización vecinal a través de capacitación y tecnología.\n" \
                  "• Proveer evidencias confiables en caso de incidentes."
    p_s2_l2.font.size = Pt(19)
    p_s2_l2.font.color.rgb = c_text_dark
    p_s2_l2.space_after = Pt(10)
    
    # Right Column: Visual Photo Card (Using real playground image!)
    insert_image_with_card(s2, 6.8, 1.8, 5.7, 4.8, "real_playground.jpg", accent_color=c_orange)
    
    add_notes(s2, "Todo gran proyecto nace de una necesidad sentida por la comunidad. En Vista al Mar, el anhelo de tener calles transitadas con tranquilidad, donde nuestros hijos puedan jugar y nuestros adultos mayores caminen sin temor, fue el motor principal. Nos dimos cuenta de que la seguridad no es solo tarea de las policías, sino que nosotros, organizados y apoyados por la tecnología, podíamos marcar una gran diferencia.")

    # ====================================================
    # SLIDE 3: La Postulación y la Larga Espera
    s3 = prs.slides.add_slide(blank_layout)
    set_bg_color(s3, c_light_warm_grey)
    add_header(s3, "Crónica de un Logro: Postulación y Espera")
    
    # 3 Column Timeline Cards
    add_card(s3, 0.75, 1.8, 3.7, 4.8, accent_color=c_text_muted)
    add_card(s3, 4.81, 1.8, 3.7, 4.8, accent_color=c_terracotta)
    add_card(s3, 8.87, 1.8, 3.7, 4.8, accent_color=c_orange)
    
    # Content Card 1
    tb_s3_c1 = s3.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(3.3), Inches(4.2))
    tf_c1 = tb_s3_c1.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "1. Postulación (2024)"
    p_c1_t.font.size = Pt(24)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = c_text_dark
    p_c1_t.space_after = Pt(10)
    p_c1_b = tf_c1.add_paragraph()
    p_c1_b.text = "En Abril de 2024, la Junta de Vecinos formaliza la postulación al fondo de Seguridad Ciudadana FNDR 8%.\n\nComienza el proceso administrativo de revisión en el Gobierno Regional."
    p_c1_b.font.size = Pt(18)
    p_c1_b.font.color.rgb = c_text_dark

    # Content Card 2
    tb_s3_c2 = s3.shapes.add_textbox(Inches(5.01), Inches(2.1), Inches(3.3), Inches(4.2))
    tf_c2 = tb_s3_c2.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "2. La Larga Espera"
    p_c2_t.font.size = Pt(24)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = c_terracotta
    p_c2_t.space_after = Pt(10)
    p_c2_b = tf_c2.add_paragraph()
    p_c2_b.text = "Retrasos gubernamentales extienden la espera hasta fines de 2025.\n\nRouter actúa como soporte técnico, explicando los plazos y manteniendo informada y unida a la comunidad."
    p_c2_b.font.size = Pt(18)
    p_c2_b.font.color.rgb = c_text_dark

    # Content Card 3
    tb_s3_c3 = s3.shapes.add_textbox(Inches(9.07), Inches(2.1), Inches(3.3), Inches(4.2))
    tf_c3 = tb_s3_c3.text_frame
    tf_c3.word_wrap = True
    p_c3_t = tf_c3.paragraphs[0]
    p_c3_t.text = "3. Unión y Triunfo"
    p_c3_t.font.size = Pt(24)
    p_c3_t.font.bold = True
    p_c3_t.font.color.rgb = c_orange
    p_c3_t.space_after = Pt(10)
    p_c3_b = tf_c3.add_paragraph()
    p_c3_b.text = "Al ver la solidez de Router, la comunidad confía y aporta documentación ágilmente ante nuevos requerimientos.\n\nFondos aprobados a fines de 2025 y transferidos en Feb 2026."
    p_c3_b.font.size = Pt(18)
    p_c3_b.font.color.rgb = c_text_dark
    
    add_notes(s3, "Postulamos en abril de 2024. Durante los largos trámites y la espera del GORE, el equipo técnico y de ingeniería de Router estuvo siempre al lado de la directiva, brindando asesoría y aclarando dudas. Esta cercanía técnica cimentó una tremenda confianza y permitió que, cuando el GORE pidió documentación de última hora a fines de 2025, el equipo la formulara en tiempo récord. Los fondos llegaron en febrero de 2026.")

    # ====================================================
    # SLIDE 4: Composición y Distribución del Fondo (Using a NATIVE PowerPoint Chart!)
    s4 = prs.slides.add_slide(blank_layout)
    set_bg_color(s4, c_light_warm_grey)
    add_header(s4, "¿Cómo se Invirtieron los Recursos del Fondo?")
    
    # Left Column: White Card with text breakdown
    add_card(s4, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s4 = s4.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s4 = tb_s4.text_frame
    tf_s4.word_wrap = True
    
    p_s4_1 = tf_s4.paragraphs[0]
    p_s4_1.text = "Distribución del Presupuesto Total ($7.000.000)"
    p_s4_1.font.size = Pt(26)
    p_s4_1.font.bold = True
    p_s4_1.font.color.rgb = c_terracotta
    p_s4_1.space_after = Pt(14)
    
    p_s4_2 = tf_s4.add_paragraph()
    p_s4_2.text = "• EQUIPAMIENTO (80.0%): $5.600.000 CLP\n" \
                 "  Adquisición e instalación de las 16 cámaras IP, red streaming e infraestructura.\n\n" \
                 "• GASTOS GENERALES (15.7%): $1.100.000 CLP\n" \
                 "  Papelería, difusión, coffee breaks para talleres y el catering/banquetes del evento de cierre.\n\n" \
                 "• RECURSOS HUMANOS (4.3%): $300.000 CLP\n" \
                 "  Honorarios profesionales de experto en seguridad y prevención de riesgos."
    p_s4_2.font.size = Pt(18)
    p_s4_2.font.color.rgb = c_text_dark
    
    # Right Column: Native editable Pie/Doughnut Chart
    add_card(s4, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    chart_data = CategoryChartData()
    chart_data.categories = ['Equipamiento (80%)', 'Gastos Generales (15.7%)', 'Recursos Humanos (4.3%)']
    chart_data.add_series('Inversión CLP', (5600000, 1100000, 300000))
    
    x, y, cx, cy = Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.0)
    chart = s4.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(14)
    chart.legend.font.name = "Arial"
    
    add_notes(s4, "Este proyecto fue financiado con un monto total de 7 millones de pesos chilenos. Siguiendo las directrices del GORE Biobío, el presupuesto se dividió de manera óptima para maximizar el impacto en seguridad. El 80%, correspondiente a $5.600.000, se destinó íntegramente a la adquisición de la tecnología de vigilancia. El porcentaje restante financió los talleres de capacitación para la comunidad, los honorarios profesionales de nuestro experto en seguridad y los elementos de difusión para que todos en el sector sepan que Vista al Mar es ahora un barrio protegido.")

    # ====================================================
    # SLIDE 5: Capacitación y Prevención (Talleres)
    s5 = prs.slides.add_slide(blank_layout)
    set_bg_color(s5, c_light_warm_grey)
    add_header(s5, "El Factor Humano: Capacitación y Prevención")
    
    # Left Card
    add_card(s5, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s5_l = s5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s5_l = tb_s5_l.text_frame
    tf_s5_l.word_wrap = True
    
    p_s5_l1 = tf_s5_l.paragraphs[0]
    p_s5_l1.text = "Charlas de Seguridad y Prevención"
    p_s5_l1.font.size = Pt(26)
    p_s5_l1.font.bold = True
    p_s5_l1.font.color.rgb = c_terracotta
    p_s5_l1.space_after = Pt(14)
    
    p_s5_l2 = tf_s5_l.add_paragraph()
    p_s5_l2.text = "• Experto Técnico a Cargo:\n" \
                  "  Ingeniero en Prevención de Riesgos o ex-miembro de las Fuerzas Armadas especialista en seguridad ($300.000 CLP en total).\n\n" \
                  "• Charlas de Capacitación Comunitaria:\n" \
                  "  Enfocadas en autocuidado vecinal, prevención del delito en espacios comunes y uso responsable del sistema.\n\n" \
                  "• Integración y Coffee Break:\n" \
                  "  Espacios de conversación para fortalecer la organización y cohesión entre vecinos."
    p_s5_l2.font.size = Pt(18)
    p_s5_l2.font.color.rgb = c_text_dark
    
    # Right Card (Photos - Using original image!)
    insert_image_with_card(s5, 6.8, 1.8, 5.7, 4.8, "vecinos_alessandri.png", accent_color=c_orange)
    
    add_notes(s5, "Para asegurar que el sistema sea realmente efectivo, se contempló un taller de capacitación por $300.000 dictado por el equipo experto de Router y un Ingeniero en Prevención de Riesgos. El taller educa a los vecinos en el autocuidado, uso ético y legal de las grabaciones, y se acompaña de un coffee break para fortalecer la organización y la cohesión vecinal.")

    # ====================================================
    # SLIDE 6: Equipamiento Tecnológico Adquirido
    s6 = prs.slides.add_slide(blank_layout)
    set_bg_color(s6, c_light_warm_grey)
    add_header(s6, "Equipamiento Tecnológico de Alta Gama")
    
    # Left: Specs Grid Card
    add_card(s6, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s6_l = s6.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s6_l = tb_s6_l.text_frame
    tf_s6_l.word_wrap = True
    
    p_s6_l1 = tf_s6_l.paragraphs[0]
    p_s6_l1.text = "Ficha Técnica: Cámaras IP Bullet"
    p_s6_l1.font.size = Pt(26)
    p_s6_l1.font.bold = True
    p_s6_l1.font.color.rgb = c_terracotta
    p_s6_l1.space_after = Pt(14)
    
    p_s6_l2 = tf_s6_l.add_paragraph()
    p_s6_l2.text = "• Alta Resolución: Imágenes de 4 Megapíxeles.\n" \
                  "• Visión Nocturna: Infrarrojo Smart Hybrid Light hasta 60m.\n" \
                  "• Micrófono: Audio ambiental de alta sensibilidad.\n" \
                  "• Enfoque Remoto: Lente varifocal motorizada.\n" \
                  "• Resistencia Climática: Certificación exterior IP67.\n" \
                  "• Calidad de Enlace: Cable de cobre UTP Categoría 5 puro."
    p_s6_l2.font.size = Pt(19)
    p_s6_l2.font.color.rgb = c_text_dark
    p_s6_l2.space_after = Pt(10)

    # Right Card: Real camera photo!
    insert_image_with_card(s6, 6.8, 1.8, 5.7, 4.8, "real_camera_red_wall.jpg", accent_color=c_orange)
    
    add_notes(s6, "Seleccionamos cámaras IP Bullet de calidad profesional con una resolución de 4 Megapíxeles, que brindan imágenes sumamente nítidas para distinguir rostros y detalles. Tienen tecnología Smart Hybrid Light con un alcance infrarrojo de hasta 60 metros en oscuridad total. Además, cuentan con micrófonos incorporados para el registro acústico, lentes varifocales motorizadas que nos permiten regular el enfoque, y una certificación IP67 que garantiza su funcionamiento sin importar el clima o la lluvia. Todo el cableado es de cobre de la más alta calidad para garantizar que la red nunca sufra caídas de rendimiento.")

    # ====================================================
    # SLIDE 7: Una Instalación Única: Tecnología Streaming
    s7 = prs.slides.add_slide(blank_layout)
    set_bg_color(s7, c_light_warm_grey)
    add_header(s7, "Instalación Innovadora Vía Streaming")
    
    # Left Big Card
    add_card(s7, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s7_l = s7.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s7_l = tb_s7_l.text_frame
    tf_s7_l.word_wrap = True
    p_s7_l1 = tf_s7_l.paragraphs[0]
    p_s7_l1.text = "Red Inalámbrica Avanzada"
    p_s7_l1.font.size = Pt(26)
    p_s7_l1.font.bold = True
    p_s7_l1.font.color.rgb = c_terracotta
    p_s7_l1.space_after = Pt(12)
    p_s7_l2 = tf_s7_l.add_paragraph()
    p_s7_l2.text = "• Enlaces Streaming de Alta Velocidad: Conectan las 16 cámaras de forma inalámbrica centralizada.\n\n" \
                  "• Cero Roturas de Calles: Se evitó romper veredas o pavimentos del barrio, eliminando costos excesivos de obra civil.\n\n" \
                  "• Precedente Regional/Nacional: Primera instalación de esta naturaleza técnica en la provincia.\n\n" \
                  "• Clima Favorable: El buen clima del sector permitió un montaje rápido y seguro en postes estratégicos."
    p_s7_l2.font.size = Pt(19)
    p_s7_l2.font.color.rgb = c_text_dark
    
    # Right Column: Visual Box (Using real generated network diagram!)
    insert_image_with_card(s7, 6.8, 1.8, 5.7, 4.8, "red_streaming.png", accent_color=c_orange)

    add_notes(s7, "Aquí viene lo más importante: este proyecto cuenta con una innovación tecnológica que marca un precedente a nivel regional e incluso nacional. En lugar de utilizar cableados tradicionales que habrían requerido romper calles y veredas, encareciendo el proyecto de manera prohibitiva, diseñamos una infraestructura de red vía streaming inalámbrico. Esto nos permitió distribuir estratégicamente 16 cámaras en diferentes puntos lejanos de la población y transmitir la señal en tiempo real sin interferencias. Además, la naturaleza nos acompañó con un clima excelente para realizar el montaje de forma rápida y segura.")

    # ====================================================
    # SLIDE 8: Mapa de Ubicación: Cobertura Vecinal
    s_map = prs.slides.add_slide(blank_layout)
    set_bg_color(s_map, c_light_warm_grey)
    add_header(s_map, "Mapa de Ubicación: Cobertura Vecinal")
    
    # Left Card
    add_card(s_map, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s_map_l = s_map.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s_map_l = tb_s_map_l.text_frame
    tf_s_map_l.word_wrap = True
    
    p_s_map_l1 = tf_s_map_l.paragraphs[0]
    p_s_map_l1.text = "Distribución Estratégica"
    p_s_map_l1.font.size = Pt(26)
    p_s_map_l1.font.bold = True
    p_s_map_l1.font.color.rgb = c_terracotta
    p_s_map_l1.space_after = Pt(14)
    
    p_s_map_l2 = tf_s_map_l.add_paragraph()
    p_s_map_l2.text = "• 16 Cámaras IP Bullet:\n" \
                      "  Distribuidas de manera uniforme para maximizar la cobertura en todos los accesos principales, pasajes internos y áreas comunes de la población Eduardo Frei Montalva.\n\n" \
                      "• Puntos Críticos de Vigilancia:\n" \
                      "  Resguardo y monitoreo continuo de vías clave como Barros Arana, Francisco Bilbao y Rigoberto Iglesias.\n\n" \
                      "• Visualización Clara:\n" \
                      "  Permite a los vecinos identificar con total precisión el campo visual y los sectores cubiertos para coordinar la prevención vecinal."
    p_s_map_l2.font.size = Pt(18)
    p_s_map_l2.font.color.rgb = c_text_dark
    
    # Right Column: Visual Box (Using the uploaded map image!)
    insert_image_with_card(s_map, 6.8, 1.8, 5.7, 4.8, "mapa_camaras.jpg", accent_color=c_orange)

    add_notes(s_map, "Este es el mapa oficial de distribución de nuestras 16 cámaras de seguridad en la población Eduardo Frei Montalva. Como pueden observar en el plano, los puntos rojos con los conos de visión verdes y amarillos representan la ubicación exacta y la dirección de enfoque de cada cámara. Diseñamos este trazado de manera sumamente estratégica en conjunto con la directiva, priorizando los pasajes interiores, los accesos principales de la población y las esquinas clave como Barros Arana y Bilbao. De esta forma, no dejamos 'puntos ciegos' significativos en los recorridos más habituales de nuestros vecinos, garantizando una cobertura comunitaria y un fuerte efecto disuasivo en toda la zona.")

    # ====================================================
    # SLIDE 9: Visualización en Vivo: Áreas Comunes
    s_live = prs.slides.add_slide(blank_layout)
    set_bg_color(s_live, c_light_warm_grey)
    add_header(s_live, "Visualización en Vivo: Monitoreo en Áreas Comunes")
    
    # Left Card
    add_card(s_live, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s_live_l = s_live.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s_live_l = tb_s_live_l.text_frame
    tf_s_live_l.word_wrap = True
    
    p_s_live_l1 = tf_s_live_l.paragraphs[0]
    p_s_live_l1.text = "Espacios Comunes Protegidos"
    p_s_live_l1.font.size = Pt(26)
    p_s_live_l1.font.bold = True
    p_s_live_l1.font.color.rgb = c_terracotta
    p_s_live_l1.space_after = Pt(14)
    
    p_s_live_l2 = tf_s_live_l.add_paragraph()
    p_s_live_l2.text = "• Sector Juegos y Multicancha:\n" \
                       "  Vigilancia constante de los principales puntos de esparcimiento familiar.\n\n" \
                       "• Calidad de Visión Nocturna:\n" \
                       "  Imágenes reales a color en penumbra e infrarrojo de gran nitidez y amplio rango.\n\n" \
                       "• Tranquilidad Comunitaria:\n" \
                       "  Resguardo continuo para que niños y jóvenes utilicen la plaza pública de forma segura."
    p_s_live_l2.font.size = Pt(19)
    p_s_live_l2.font.color.rgb = c_text_dark
    
    # Right Column: Visual Box (Using the uploaded park screenshot!)
    insert_image_with_card(s_live, 6.8, 1.8, 5.7, 4.8, "vista_parque.jpg", accent_color=c_orange)

    add_notes(s_live, "Como una imagen vale más que mil palabras, queremos mostrarles de forma muy concreta cómo se ven nuestras cámaras instaladas y operativas en plena noche. Esta captura que ven en pantalla corresponde al sector del parque, enfocando la zona de juegos infantiles y la multicancha. Como pueden notar, gracias a la tecnología de iluminación híbrida y la resolución de 4 Megapíxeles, la imagen nocturna es sumamente nítida, brillante y con un rango de visión muy amplio. De esta manera, garantizamos que los espacios públicos donde se reúnen nuestros niños y jóvenes estén constantemente monitoreados, trayendo tranquilidad a todas las familias que viven cerca del parque.")

    # ====================================================
    # SLIDE 10: Claridad Tecnológica: ¿Por qué No Leen Patentes?
    s8 = prs.slides.add_slide(blank_layout)
    set_bg_color(s8, c_light_warm_grey)
    add_header(s8, "Transparencia: Límites Técnicos del Equipamiento")
    
    # 2 Column Cost comparison layout
    add_card(s8, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    add_card(s8, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    
    # Column 1: Our cameras
    tb_s8_c1 = s8.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s8_c1 = tb_s8_c1.text_frame
    tf_s8_c1.word_wrap = True
    tf_s8_c1.paragraphs[0].text = "Cámaras 4 Megapíxeles Adquiridas"
    tf_s8_c1.paragraphs[0].font.bold = True
    tf_s8_c1.paragraphs[0].font.size = Pt(24)
    tf_s8_c1.paragraphs[0].font.color.rgb = c_terracotta
    tf_s8_c1.paragraphs[0].space_after = Pt(14)
    p_s8_c1_b = tf_s8_c1.add_paragraph()
    p_s8_c1_b.text = "• Cobertura: 16 puntos estratégicos protegidos.\n\n" \
                    "• Calidad: Videos de alta resolución que registran la escena general, vehículos, colores, modelos y rostros de personas.\n\n" \
                    "• Decisión: Priorizar la extensión territorial del sistema de vigilancia para el resguardo de todos los pasajes."
    p_s8_c1_b.font.size = Pt(19)
    p_s8_c1_b.font.color.rgb = c_text_dark

    # Column 2: LPR cameras
    tb_s8_c2 = s8.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s8_c2 = tb_s8_c2.text_frame
    tf_s8_c2.word_wrap = True
    tf_s8_c2.paragraphs[0].text = "Lector de Patentes (Tecnología LPR)"
    tf_s8_c2.paragraphs[0].font.bold = True
    tf_s8_c2.paragraphs[0].font.size = Pt(24)
    tf_s8_c2.paragraphs[0].font.color.rgb = c_orange
    tf_s8_c2.paragraphs[0].space_after = Pt(14)
    p_s8_c2_b = tf_s8_c2.add_paragraph()
    p_s8_c2_b.text = "• Limitación: NO incluido en este equipamiento.\n\n" \
                    "• Factor de Costo: Un solo sistema LPR real y homologado de lectura automatizada cuesta 10 veces más que una cámara estándar.\n\n" \
                    "• Impacto Presupuestario: Elegir LPR habría limitado la protección a una sola esquina de la población, dejando al resto de la comunidad sin resguardo."
    p_s8_c2_b.font.size = Pt(19)
    p_s8_c2_b.font.color.rgb = c_text_dark
    
    add_notes(s8, "Queremos hacer una aclaración técnica muy relevante para evitar falsas expectativas y ser 100% transparentes: estas cámaras NO leen patentes automáticas. ¿Por qué se tomó esta decisión? Un sistema lector de patentes real y homologado (tecnología LPR) cuesta hasta 10 veces más que una cámara de seguridad de alta calidad. Si hubiésemos optado por lectores de patentes, el presupuesto solo nos habría alcanzado para vigilar una o dos esquinas de la población, dejando al resto de los vecinos desprotegidos. Al elegir cámaras de alta resolución estándar de 4MP, pudimos proteger 16 puntos estratégicos. Si ocurre un incidente, igual podemos ver el auto, su color, modelo y dirección, pero la lectura automática de la placa de un vehículo en movimiento requiere otra gama de costos no viable en este tipo de fondos comunitarios.")

    # ====================================================
    # SLIDE 11: Grabación y Almacenamiento Optimizado
    s9 = prs.slides.add_slide(blank_layout)
    set_bg_color(s9, c_light_warm_grey)
    add_header(s9, "Estrategia de Grabación para Máxima Utilidad")
    
    # Two Columns: Left text, Right Graphic Box
    add_card(s9, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s9_l = s9.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s9_l = tb_s9_l.text_frame
    tf_s9_l.word_wrap = True
    
    p_s9_l1 = tf_s9_l.paragraphs[0]
    p_s9_l1.text = "Calidad sobre Cantidad de Días"
    p_s9_l1.font.size = Pt(26)
    p_s9_l1.font.bold = True
    p_s9_l1.font.color.rgb = c_terracotta
    p_s9_l1.space_after = Pt(14)
    
    p_s9_l2 = tf_s9_l.add_paragraph()
    p_s9_l2.text = "• Capacidad Técnica: Cada cámara cuenta con memoria SD de 64GB y puede grabar hasta 15 días continuos en baja calidad.\n\n" \
                  "• Ajuste Adoptado: El almacenamiento central está configurado exactamente para 5 días continuos de video.\n\n" \
                  "• Justificación de Router: Al concentrar los 64GB en 5 días, podemos transmitir y grabar en la máxima nitidez y fotogramas por segundo (FPS). Esto permite que el video sirva como evidencia útil de rostros e incidentes."
    p_s9_l2.font.size = Pt(19)
    p_s9_l2.font.color.rgb = c_text_dark

    # Right Card: Native Chart showing 15 Days vs 5 Days resolution quality index
    add_card(s9, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    chart_data_s9 = CategoryChartData()
    chart_data_s9.categories = ['Configurado 5 días', 'Configurado 15 días']
    chart_data_s9.add_series('Índice de Nitidez (1-100)', (95, 30))
    chart_data_s9.add_series('FPS (Fotogramas x seg)', (25, 8))
    
    x_s9, y_s9, cx_s9, cy_s9 = Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.0)
    chart_s9 = s9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x_s9, y_s9, cx_s9, cy_s9, chart_data_s9
    ).chart
    chart_s9.has_legend = True
    chart_s9.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart_s9.legend.font.size = Pt(14)
    
    add_notes(s9, "¿Cómo funciona el almacenamiento de las grabaciones? Cada cámara cuenta con almacenamiento interno en tarjetas de memoria industrial de 64 Gigabytes. Técnicamente, las cámaras podrían grabar hasta 15 días en una calidad comprimida. Sin embargo, para que el sistema sea realmente útil, tomamos la decisión estratégica de configurarlas para grabar los últimos 5 días. Al acortar el rango, podemos configurarlas en la máxima resolución y calidad posible de imagen. Es mejor tener 5 días de video nítido e impecable donde se distingan bien las caras y detalles, que tener 15 días borrosos o pixelados que no sirvan como evidencia en tribunales. Las pruebas fueron exitosas y todo el sistema está rindiendo conforme a lo comprometido.")

    # ====================================================
    # SLIDE 12: Responsabilidad Local: Encargado de Cámaras
    s10 = prs.slides.add_slide(blank_layout)
    set_bg_color(s10, c_light_warm_grey)
    add_header(s10, "Custodia Vecinal y Resguardo del Sistema")
    
    # Left Column: Info Card
    add_card(s10, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s10_l = s10.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s10_l = tb_s10_l.text_frame
    tf_s10_l.word_wrap = True
    
    p_s10_l1 = tf_s10_l.paragraphs[0]
    p_s10_l1.text = "Custodia y Nexo Vecinal"
    p_s10_l1.font.size = Pt(26)
    p_s10_l1.font.bold = True
    p_s10_l1.font.color.rgb = c_terracotta
    p_s10_l1.space_after = Pt(14)
    
    p_s10_l2 = tf_s10_l.add_paragraph()
    p_s10_l2.text = "• Encargado Designado: Encargado de Seguridad Vecinal designado por la directiva.\n\n" \
                   "• Funciones Clave:\n" \
                   "  - Custodiar físicamente y resguardar la central de monitoreo.\n" \
                   "  - Servir de nexo técnico oficial con la Directiva.\n" \
                   "  - Resguardar bajo estricta confidencialidad los datos y grabaciones para proteger la privacidad de los vecinos."
    p_s10_l2.font.size = Pt(19)
    p_s10_l2.font.color.rgb = c_text_dark
    
    # Right Column: Photo Card (Using real photo of custodian and neighbors!)
    insert_image_with_card(s10, 6.8, 1.8, 5.7, 4.8, "real_neighbors.jpg", accent_color=c_orange)

    add_notes(s10, "Un sistema de seguridad comunitario requiere una persona de absoluta confianza a cargo de los equipos para garantizar la privacidad y el uso responsable. La asamblea y la directiva han asignado formalmente esta importante labor a nuestro un vecino de absoluta confianza designado formalmente. Este custodio será el custodio técnico de la central de monitoreo, velando por que los equipos funcionen correctamente y actuando bajo un estricto principio de confidencialidad y respeto a la privacidad de todas las familias.")

    # ====================================================
    # SLIDE 13: Protocolo de Seguridad y Privacidad
    s11 = prs.slides.add_slide(blank_layout)
    set_bg_color(s11, c_light_warm_grey)
    add_header(s11, "Protocolo de Seguridad y Privacidad")
    
    # Left Card
    add_card(s11, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s11_l = s11.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s11_l = tb_s11_l.text_frame
    tf_s11_l.word_wrap = True
    p_s11_l1 = tf_s11_l.paragraphs[0]
    p_s11_l1.text = "¿Cuándo se Activa el Protocolo?"
    p_s11_l1.font.size = Pt(26)
    p_s11_l1.font.bold = True
    p_s11_l1.font.color.rgb = c_terracotta
    p_s11_l1.space_after = Pt(14)
    p_s11_l2 = tf_s11_l.add_paragraph()
    p_s11_l2.text = "• Eventos Delictuales: Sospecha fundada o comisión de delitos (robos, asaltos, hurtos).\n\n" \
                   "• Daños Materiales: Daños o vandalismo en bienes comunes o propiedad de los vecinos.\n\n" \
                   "• Accidentes Viales: Choques o atropellos dentro del sector de la población."
    p_s11_l2.font.size = Pt(19)
    p_s11_l2.font.color.rgb = c_text_dark

    # Right Card
    add_card(s11, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    tb_s11_r = s11.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s11_r = tb_s11_r.text_frame
    tf_s11_r.word_wrap = True
    p_s11_r1 = tf_s11_r.paragraphs[0]
    p_s11_r1.text = "Límites Técnicos y de Privacidad"
    p_s11_r1.font.size = Pt(26)
    p_s11_r1.font.bold = True
    p_s11_r1.font.color.rgb = c_orange
    p_s11_r1.space_after = Pt(14)
    p_s11_r2 = tf_s11_r.add_paragraph()
    p_s11_r2.text = "• Espacio Público: Orientación exclusiva de las lentes hacia calles, veredas y pasajes públicos.\n\n" \
                   "• Privacidad Hogareña: Calibración estricta para evitar la captación de patios interiores y ventanas de los vecinos.\n\n" \
                   "• Uso Restringido: Queda estrictamente prohibido utilizar grabaciones para fines de esparcimiento o no autorizados."
    p_s11_r2.font.size = Pt(19)
    p_s11_r2.font.color.rgb = c_text_dark

    add_notes(s11, "Para resguardar el derecho a la privacidad de todos los hogares, hemos diseñado un protocolo formal y transparente de acceso a las grabaciones. Las cámaras no son para 'espiar' el día a día del barrio; son para proteger. Por lo tanto, si un vecino necesita revisar una grabación debido a un robo, choque, daño o cualquier ilícito, no debe acudir directamente al Custodio Técnico. El canal correcto es presentar una solicitud formal ante la Directiva de la Junta de Vecinos. La directiva analizará el caso de forma reservada, autorizará la extracción del video de la fecha y hora indicadas, y el Custodio procesará el clip para que sea enviado de manera rápida y segura vía WhatsApp al vecino o directamente a las policías como medio de prueba para la denuncia formal.")

    # ====================================================
    # SLIDE 14: Pasos Formales de Solicitud (Protocolo)
    s12 = prs.slides.add_slide(blank_layout)
    set_bg_color(s12, c_light_warm_grey)
    add_header(s12, "Pasos Formales de Solicitud (Protocolo)")
    
    # 4 Grid Steps
    add_card(s12, 0.75, 1.8, 5.7, 2.3, accent_color=c_terracotta)
    add_card(s12, 6.8, 1.8, 5.7, 2.3, accent_color=c_orange)
    add_card(s12, 0.75, 4.3, 5.7, 2.3, accent_color=c_text_muted)
    add_card(s12, 6.8, 4.3, 5.7, 2.3, accent_color=c_dark_charcoal)
    
    # Step 1
    tb_s12_s1 = s12.shapes.add_textbox(Inches(0.95), Inches(2.05), Inches(5.3), Inches(1.9))
    tf_s12_s1 = tb_s12_s1.text_frame
    tf_s12_s1.word_wrap = True
    tf_s12_s1.paragraphs[0].text = "Paso 1: Solicitud formal a la Directiva"
    tf_s12_s1.paragraphs[0].font.bold = True
    tf_s12_s1.paragraphs[0].font.size = Pt(20)
    tf_s12_s1.paragraphs[0].font.color.rgb = c_terracotta
    p_s12_s1_b = tf_s12_s1.add_paragraph()
    p_s12_s1_b.text = "El vecino afectado formaliza su solicitud a la directiva por escrito o a través del WhatsApp institucional. No debe acudir directamente al cuidador técnico."
    p_s12_s1_b.font.size = Pt(16)
    p_s12_s1_b.font.color.rgb = c_text_dark
    p_s12_s1_b.space_before = Pt(4)

    # Step 2
    tb_s12_s2 = s12.shapes.add_textbox(Inches(7.0), Inches(2.05), Inches(5.3), Inches(1.9))
    tf_s12_s2 = tb_s12_s2.text_frame
    tf_s12_s2.word_wrap = True
    tf_s12_s2.paragraphs[0].text = "Paso 2: Datos Técnicos Requeridos"
    tf_s12_s2.paragraphs[0].font.bold = True
    tf_s12_s2.paragraphs[0].font.size = Pt(20)
    tf_s12_s2.paragraphs[0].font.color.rgb = c_orange
    p_s12_s2_b = tf_s12_s2.add_paragraph()
    p_s12_s2_b.text = "Especificar detalladamente: fecha del incidente, rango de hora aproximado y ubicación exacta o cámara más cercana al suceso."
    p_s12_s2_b.font.size = Pt(16)
    p_s12_s2_b.font.color.rgb = c_text_dark
    p_s12_s2_b.space_before = Pt(4)

    # Step 3
    tb_s12_s3 = s12.shapes.add_textbox(Inches(0.95), Inches(4.55), Inches(5.3), Inches(1.9))
    tf_s12_s3 = tb_s12_s3.text_frame
    tf_s12_s3.word_wrap = True
    tf_s12_s3.paragraphs[0].text = "Paso 3: Extracción y Envío por WhatsApp"
    tf_s12_s3.paragraphs[0].font.bold = True
    tf_s12_s3.paragraphs[0].font.size = Pt(20)
    tf_s12_s3.paragraphs[0].font.color.rgb = c_text_muted
    p_s12_s3_b = tf_s12_s3.add_paragraph()
    p_s12_s3_b.text = "La directiva aprueba y canaliza la solicitud con el Custodio Técnico, quien extrae el fragmento de video y lo envía de forma segura vía WhatsApp."
    p_s12_s3_b.font.size = Pt(16)
    p_s12_s3_b.font.color.rgb = c_text_dark
    p_s12_s3_b.space_before = Pt(4)

    # Step 4
    tb_s12_s4 = s12.shapes.add_textbox(Inches(7.0), Inches(4.55), Inches(5.3), Inches(1.9))
    tf_s12_s4 = tb_s12_s4.text_frame
    tf_s12_s4.word_wrap = True
    tf_s12_s4.paragraphs[0].text = "Límite Crítico: Plazo de 5 Días"
    tf_s12_s4.paragraphs[0].font.bold = True
    tf_s12_s4.paragraphs[0].font.size = Pt(20)
    tf_s12_s4.paragraphs[0].font.color.rgb = c_dark_charcoal
    p_s12_s4_b = tf_s12_s4.add_paragraph()
    p_s12_s4_b.text = "IMPORTANTE: Realizar el requerimiento antes de cumplirse los 5 días del hecho. Posterior a esto, el almacenamiento se sobrescribirá."
    p_s12_s4_b.font.size = Pt(16)
    p_s12_s4_b.font.color.rgb = c_text_dark
    p_s12_s4_b.space_before = Pt(4)

    add_notes(s12, "Este es un detalle fundamental que debemos memorizar: el plazo máximo para solicitar una grabación es de 5 días. Debido a que el sistema sobrescribe las grabaciones antiguas para mantener la mejor calidad del video, si un incidente ocurre un lunes, la directiva debe recibir la solicitud antes del viernes de esa misma semana. La solicitud debe indicar con la mayor precisión posible el día, la hora y el lugar del hecho. Con esto garantizamos un proceso ordenado, rápido y sobre todo legal ante cualquier requerimiento de los tribunales.")

    # ====================================================
    # SLIDE 15: Aclarando Dudas Comunitarias (FAQ)
    s13 = prs.slides.add_slide(blank_layout)
    set_bg_color(s13, c_light_warm_grey)
    add_header(s13, "Aclarando Dudas Frecuentes (FAQ)")
    
    # 2 columns of FAQ cards
    add_card(s13, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    add_card(s13, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    
    # FAQ Column 1
    tb_s13_c1 = s13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_s13_c1 = tb_s13_c1.text_frame
    tf_s13_c1.word_wrap = True
    
    p_faq_1 = tf_s13_c1.paragraphs[0]
    p_faq_1.text = "¿Quién financia la electricidad?"
    p_faq_1.font.bold = True
    p_faq_1.font.size = Pt(22)
    p_faq_1.font.color.rgb = c_terracotta
    p_faq_1.space_after = Pt(2)
    p_faq_1_a = tf_s13_c1.add_paragraph()
    p_faq_1_a.text = "Costo mensual de $1.200 pesos de electricidad por cámara, asumido por los vecinos que apoyaron en el proyecto (quienes facilitaron la conexión eléctrica)."
    p_faq_1_a.font.size = Pt(17)
    p_faq_1_a.font.color.rgb = c_text_dark
    p_faq_1_a.space_after = Pt(10)
    
    p_faq_2 = tf_s13_c1.add_paragraph()
    p_faq_2.text = "¿Tiene costo mensual para el vecino?"
    p_faq_2.font.bold = True
    p_faq_2.font.size = Pt(22)
    p_faq_2.font.color.rgb = c_terracotta
    p_faq_2.space_after = Pt(2)
    p_faq_2_a = tf_s13_c1.add_paragraph()
    p_faq_2_a.text = "No. El fondo FNDR financió la instalación y Router entrega el equipamiento cerrado y libre de cuotas fijas mensuales ($0 CLP de mantención fija)."
    p_faq_2_a.font.size = Pt(17)
    p_faq_2_a.font.color.rgb = c_text_dark

    # FAQ Column 2
    tb_s13_c2 = s13.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_s13_c2 = tb_s13_c2.text_frame
    tf_s13_c2.word_wrap = True
    
    p_faq_3 = tf_s13_c2.paragraphs[0]
    p_faq_3.text = "¿Se transmiten por Internet pública?"
    p_faq_3.font.bold = True
    p_faq_3.font.size = Pt(22)
    p_faq_3.font.color.rgb = c_orange
    p_faq_3.space_after = Pt(2)
    p_faq_3_a = tf_s13_c2.add_paragraph()
    p_faq_3_a.text = "No. El sistema utiliza una red local de streaming inalámbrico propiedad de la Junta de Vecinos. No se suben clips a redes abiertas o nubes públicas."
    p_faq_3_a.font.size = Pt(17)
    p_faq_3_a.font.color.rgb = c_text_dark
    p_faq_3_a.space_after = Pt(10)
    
    p_faq_4 = tf_s13_c2.add_paragraph()
    p_faq_4.text = "¿Cada cuánto una mantención?"
    p_faq_4.font.bold = True
    p_faq_4.font.size = Pt(22)
    p_faq_4.font.color.rgb = c_orange
    p_faq_4.space_after = Pt(2)
    p_faq_4_a = tf_s13_c2.add_paragraph()
    p_faq_4_a.text = "Las mantenciones debieran ser idealmente cada 3 meses para limpiar las cámaras y ver el correcto funcionamiento."
    p_faq_4_a.font.size = Pt(17)
    p_faq_4_a.font.color.rgb = c_text_dark

    add_notes(s13, "Es muy común que en estas reuniones surjan dudas. Por ejemplo, sobre el costo eléctrico: es de unos $1.200 pesos mensuales por cámara, costo que es asumido generosamente por los vecinos colaboradores que prestan la energía de sus hogares. Tampoco hay cuotas mensuales para los vecinos; el sistema es propio de la Junta. Además, las grabaciones viajan de forma segura por una red inalámbrica local y no se suben a internet. Para que el sistema dure en el tiempo, se recomienda realizar mantenciones preventivas cada 3 meses.")

    # ====================================================
    # SLIDE 16: Garantía Oficial del Sistema
    s14 = prs.slides.add_slide(blank_layout)
    set_bg_color(s14, c_light_warm_grey)
    add_header(s14, "Garantía Oficial del Sistema")
    
    # 2 columns of Warranty cards
    add_card(s14, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    add_card(s14, 6.8, 1.8, 5.7, 4.8, accent_color=c_orange)
    
    # Column 1: Cobertura
    tb_s14_c1 = s14.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s14_c1 = tb_s14_c1.text_frame
    tf_s14_c1.word_wrap = True
    
    p_w1 = tf_s14_c1.paragraphs[0]
    p_w1.text = "Cobertura de Garantía Router"
    p_w1.font.bold = True
    p_w1.font.size = Pt(26)
    p_w1.font.color.rgb = c_terracotta
    p_w1.space_after = Pt(14)
    
    p_w1_b = tf_s14_c1.add_paragraph()
    p_w1_b.text = "• Plazo de Cobertura: 6 meses de garantía a partir del primero de mayo.\n\n" \
                  "• Alcance Técnico: Contempla cualquier falla técnica del producto o defectos en la instalación del equipamiento realizada por nuestra empresa."
    p_w1_b.font.size = Pt(19)
    p_w1_b.font.color.rgb = c_text_dark
    
    # Column 2: Exclusiones
    tb_s14_c2 = s14.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.2), Inches(4.2))
    tf_s14_c2 = tb_s14_c2.text_frame
    tf_s14_c2.word_wrap = True
    
    p_w2 = tf_s14_c2.paragraphs[0]
    p_w2.text = "Exclusiones de Cobertura"
    p_w2.font.bold = True
    p_w2.font.size = Pt(26)
    p_w2.font.color.rgb = c_orange
    p_w2.space_after = Pt(14)
    
    p_w2_b = tf_s14_c2.add_paragraph()
    p_w2_b.text = "• No Cubierto: No se considera bajo garantía la manipulación indebida, intervenciones de terceros o modificaciones físicas no autorizadas.\n\n" \
                  "• Siniestros Excluidos: Daños provocados por vandalismo, robos, sobretensiones eléctricas externas u otros factores climáticos destructivos extremos."
    p_w2_b.font.size = Pt(19)
    p_w2_b.font.color.rgb = c_text_dark
    
    add_notes(s14, "Nuestra empresa Router entrega una garantía técnica de 6 meses desde la recepción de la obra, respaldada por un equipo de soporte de respuesta rápida. Esto cubre desperfectos de fábrica o montaje. La garantía no cubre manipulación no autorizada, robos o vandalismo, por lo que es vital cuidar los equipos en comunidad.")

    # ====================================================
    # SLIDE 17: Espacio de Consultas y Diálogo
    s15 = prs.slides.add_slide(blank_layout)
    set_bg_color(s15, c_light_warm_grey)
    add_header(s15, "Foro Abierto: Preguntas y Diálogo Vecinal")
    
    add_card(s15, 0.75, 1.8, 11.83, 4.8, accent_color=c_terracotta)
    tb_s15 = s15.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf_s15 = tb_s15.text_frame
    tf_s15.word_wrap = True
    
    p_s15_1 = tf_s15.paragraphs[0]
    p_s15_1.text = "Espacio de consultas para la asamblea vecinal"
    p_s15_1.font.size = Pt(34)
    p_s15_1.font.bold = True
    p_s15_1.font.color.rgb = c_terracotta
    p_s15_1.alignment = PP_ALIGN.CENTER
    p_s15_1.space_after = Pt(24)
    
    p_s15_2 = tf_s15.add_paragraph()
    p_s15_2.text = "Tu opinión y tus dudas son muy importantes.\n\n" \
                  "Por favor, levanta la mano y comparte tus inquietudes con la Directiva de la Junta de Vecinos Eduardo Frei Montalva y el Equipo Técnico de Router."
    p_s15_2.font.size = Pt(24)
    p_s15_2.font.color.rgb = c_text_dark
    p_s15_2.alignment = PP_ALIGN.CENTER
    
    add_notes(s15, "Ahora abrimos la palabra para ustedes. Queremos escuchar sus consultas, sugerencias o impresiones. Ninguna pregunta es menos importante; al contrario, este es el momento para que todos nos vayamos a casa con total claridad sobre el funcionamiento de nuestra nueva red de seguridad comunitaria.")

    # ====================================================
    # SLIDE 18: Difusión y Video Promocional
    s16 = prs.slides.add_slide(blank_layout)
    set_bg_color(s16, c_light_warm_grey)
    add_header(s16, "¡Vista al Mar Unida en Redes Sociales!")
    
    # Left Card
    add_card(s16, 0.75, 1.8, 5.7, 4.8, accent_color=c_terracotta)
    tb_s16_l = s16.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf_s16_l = tb_s16_l.text_frame
    tf_s16_l.word_wrap = True
    
    p_s16_l1 = tf_s16_l.paragraphs[0]
    p_s16_l1.text = "La Seguridad se Construye Comunicando"
    p_s16_l1.font.size = Pt(26)
    p_s16_l1.font.bold = True
    p_s16_l1.font.color.rgb = c_terracotta
    p_s16_l1.space_after = Pt(12)
    
    p_s16_l2 = tf_s16_l.add_paragraph()
    p_s16_l2.text = "• ¡Síguenos en Redes Sociales! Sigue a Router y mantente informado de consejos y proyectos de seguridad.\n\n" \
                   "• Foco Disuasivo: Comparte el video promocional para alertar a extraños que el barrio está protegido.\n\n" \
                   "• Viraliza el Proyecto: Comparte el video y actas en tus grupos de pasajes y WhatsApp.\n\n" \
                   "Nuestras Redes Oficiales:\n" \
                   "  - Instagram: @routersolucionestecnologicas\n" \
                   "  - Facebook: /routersolucionestecnologicas"
    p_s16_l2.font.size = Pt(16)
    p_s16_l2.font.color.rgb = c_text_dark
    
    # Right Card (Real playground aerial drone image!)
    insert_image_with_card(s16, 6.8, 1.8, 5.7, 4.8, "real_playground_aerial.jpg", accent_color=c_orange)
    
    add_notes(s16, "Parte integral del proyecto financiado por el GORE es la difusión. Creamos un video promocional de alta calidad que muestra el trabajo realizado, la instalación y la unión de nuestros vecinos. Les pedimos que lo compartan activamente en sus redes sociales y grupos de WhatsApp vecinales. Mientras más personas sepan que en Vista al Mar estamos organizados, coordinados y monitoreados con 16 cámaras de alta tecnología, menos delincuentes se atreverán a ingresar a nuestras calles.")

    # ====================================================
    # SLIDE 19: Promoción Especial Vecinos [NEW]
    s_promo = prs.slides.add_slide(blank_layout)
    set_bg_color(s_promo, c_light_warm_grey)
    add_header(s_promo, "Promoción Especial para tu Hogar", category="BENEFICIO COMUNITARIO")
    
    # Left Card
    add_card(s_promo, 0.75, 1.8, 5.7, 4.8, accent_color=c_orange)
    tb_sp_l = s_promo.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf_sp_l = tb_sp_l.text_frame
    tf_sp_l.word_wrap = True
    
    p_sp_l1 = tf_sp_l.paragraphs[0]
    p_sp_l1.text = "Instalación de 1 Cámara Domiciliaria"
    p_sp_l1.font.size = Pt(26)
    p_sp_l1.font.bold = True
    p_sp_l1.font.color.rgb = c_orange
    p_sp_l1.space_after = Pt(6)
    
    p_sp_l2 = tf_sp_l.add_paragraph()
    p_sp_l2.text = "Precio Preferencial: $90.000 CLP (Pago Único)"
    p_sp_l2.font.size = Pt(19)
    p_sp_l2.font.bold = True
    p_sp_l2.font.color.rgb = c_terracotta
    p_sp_l2.space_after = Pt(10)
    
    p_sp_l3 = tf_sp_l.add_paragraph()
    p_sp_l3.text = "• Cámara de Alta Gama: Full HD, visión nocturna, audio.\n" \
                   "• Almacenamiento Local: Tarjeta Micro SD de 64GB incluida.\n" \
                   "• Todo Incluido: Instalación técnica, configuración y garantía Router.\n" \
                   "• Vigencia Limitada: Promoción válida solo hasta el 15 de Julio de 2026."
    p_sp_l3.font.size = Pt(15)
    p_sp_l3.font.color.rgb = c_text_dark
    p_sp_l3.space_after = Pt(12)
    
    # WhatsApp Highlight Box in Left Card
    promo_shape = s_promo.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.2), Inches(5.4), Inches(1.2))
    promo_shape.fill.solid()
    promo_shape.fill.fore_color.rgb = RGBColor(240, 253, 244)  # Light green tint
    promo_shape.line.color.rgb = RGBColor(22, 163, 74)
    promo_shape.line.width = Pt(1.5)
    
    tf_p = promo_shape.text_frame
    tf_p.word_wrap = True
    tf_p.margin_top = Inches(0.08)
    tf_p.margin_bottom = Inches(0.08)
    tf_p.margin_left = Inches(0.15)
    tf_p.margin_right = Inches(0.15)
    
    p1 = tf_p.paragraphs[0]
    p1.text = "¡RESERVA TU INSTALACIÓN HOY!"
    p1.font.bold = True
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(22, 163, 74)
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(2)
    
    p2 = tf_p.add_paragraph()
    p2.text = "WhatsApp: +56 9 7858 9090"
    p2.font.bold = True
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(22, 163, 74)
    p2.alignment = PP_ALIGN.CENTER
    
    # Right Card: camera_alessandri.png
    insert_image_with_card(s_promo, 6.8, 1.8, 5.7, 4.8, "camera_alessandri.png", accent_color=c_orange)
    
    add_notes(s_promo, "Para apoyar a toda la comunidad de la Junta de Vecinos de Vista al Mar en su seguridad familiar, desde Router hemos preparado un beneficio exclusivo para ustedes. Se trata de una promoción especial para instalar una cámara particular en sus hogares por un valor preferencial de $90.000 pesos. Este kit contempla una cámara de alta gama Full HD, equipada con audio y una tarjeta micro SD de 64GB para almacenar los videos. Es muy importante destacar que este precio incluye tanto la instalación técnica como la garantía oficial de Router. Tengan en consideración que la vigencia de esta promoción es limitada y vencerá de forma impostergable el 15 de julio de 2026. Si desean coordinar, pueden agendar desde ya al WhatsApp que aparece en pantalla: +56 9 7858 9090.")

    # ====================================================
    # SLIDE 20: Conclusión y Agradecimientos (Dark charcoal background)
    s17 = prs.slides.add_slide(blank_layout)
    set_bg_color(s17, c_dark_charcoal)
    
    # Left Terracotta Stripe
    block_s17 = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    block_s17.fill.solid()
    block_s17.fill.fore_color.rgb = c_terracotta
    block_s17.line.fill.background()
    
    # Title & Subtitle Box
    tb_s17_top = s17.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(3.0))
    tf_s17_top = tb_s17_top.text_frame
    tf_s17_top.word_wrap = True
    tf_s17_top.margin_left = tf_s17_top.margin_top = tf_s17_top.margin_right = tf_s17_top.margin_bottom = 0
    
    p_s17_1 = tf_s17_top.paragraphs[0]
    p_s17_1.text = "¡JUNTOS HACEMOS UN BARRIO MÁS SEGURO!"
    p_s17_1.font.name = "Arial"
    p_s17_1.font.size = Pt(46)
    p_s17_1.font.bold = True
    p_s17_1.font.color.rgb = c_white
    p_s17_1.space_after = Pt(16)
    
    p_s17_2 = tf_s17_top.add_paragraph()
    p_s17_2.text = '"La tecnología vigila, pero la comunidad organizada protege."'
    p_s17_2.font.name = "Arial"
    p_s17_2.font.size = Pt(26)
    p_s17_2.font.bold = True
    p_s17_2.font.color.rgb = c_orange
    
    # Footer separator line
    line_s17 = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.2), Inches(11.0), Inches(0.02))
    line_s17.fill.solid()
    line_s17.fill.fore_color.rgb = RGBColor(71, 85, 105)
    line_s17.line.fill.background()
    
    # Left Footer Box: Agradecimientos
    tb_s17_agrad = s17.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(5.5), Inches(2.2))
    tf_s17_agrad = tb_s17_agrad.text_frame
    tf_s17_agrad.word_wrap = True
    tf_s17_agrad.margin_left = tf_s17_agrad.margin_top = tf_s17_agrad.margin_right = tf_s17_agrad.margin_bottom = 0
    
    p_agrad_title = tf_s17_agrad.paragraphs[0]
    p_agrad_title.text = "AGRADECIMIENTOS ESPECIALES:"
    p_agrad_title.font.name = "Arial"
    p_agrad_title.font.bold = True
    p_agrad_title.font.size = Pt(16)
    p_agrad_title.font.color.rgb = RGBColor(148, 163, 184)
    p_agrad_title.space_after = Pt(8)
    
    p_agrad_body = tf_s17_agrad.add_paragraph()
    p_agrad_body.text = "• Gobierno Regional del Biobío (GORE Biobío)\n" \
                        "• Directiva Junta de Vecinos Población Eduardo Frei Montalva\n" \
                        "• A todos los Vecinos y Vecinas de Vista al Mar"
    p_agrad_body.font.name = "Arial"
    p_agrad_body.font.size = Pt(16)
    p_agrad_body.font.color.rgb = c_text_muted
    
    # Vertical Separator between footers
    line_vert = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(4.5), Inches(0.02), Inches(1.8))
    line_vert.fill.solid()
    line_vert.fill.fore_color.rgb = RGBColor(71, 85, 105)
    line_vert.line.fill.background()
    
    # Right Footer Box: Contacto Router
    tb_s17_contact = s17.shapes.add_textbox(Inches(7.3), Inches(4.5), Inches(5.0), Inches(2.2))
    tf_s17_contact = tb_s17_contact.text_frame
    tf_s17_contact.word_wrap = True
    tf_s17_contact.margin_left = tf_s17_contact.margin_top = tf_s17_contact.margin_right = tf_s17_contact.margin_bottom = 0
    
    p_contact_title = tf_s17_contact.paragraphs[0]
    p_contact_title.text = "SOCIO TECNOLÓGICO:"
    p_contact_title.font.name = "Arial"
    p_contact_title.font.bold = True
    p_contact_title.font.size = Pt(16)
    p_contact_title.font.color.rgb = RGBColor(148, 163, 184)
    p_contact_title.space_after = Pt(8)
    
    p_contact_phone = tf_s17_contact.add_paragraph()
    p_contact_phone.text = "WhatsApp: +56 9 7858 9090"
    p_contact_phone.font.name = "Arial"
    p_contact_phone.font.bold = True
    p_contact_phone.font.size = Pt(22)
    p_contact_phone.font.color.rgb = RGBColor(22, 163, 74)
    p_contact_phone.space_after = Pt(4)
    
    p_contact_sub = tf_s17_contact.add_paragraph()
    p_contact_sub.text = "Router • Soluciones Tecnológicas y de Ingeniería"
    p_contact_sub.font.name = "Arial"
    p_contact_sub.font.size = Pt(16)
    p_contact_sub.font.color.rgb = c_text_muted
    
    add_notes(s17, "Llegamos al final de esta presentación, pero esto es solo el comienzo de una etapa más segura para todos. Queremos reiterar nuestro más profundo agradecimiento al Gobierno Regional del Biobío por confiar en nosotros y proveer los recursos; a la directiva que no durmió recopilando papeles y empujando el proyecto; y a cada uno de ustedes por su paciencia durante la larga espera y su apoyo activo en la instalación. Como dice nuestro lema final: las cámaras vigilan, pero es la organización vecinal y la solidaridad entre nosotros lo que realmente protege a nuestras familias. ¡Muchas gracias a todos y a disfrutar de un barrio más seguro!")

    # Save presentation
    output_path = "/Users/ricardomarimodinger/.gemini/antigravity/scratch/ricardo-ai-system/presentacion_vista_al_mar.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_premium_presentation()
